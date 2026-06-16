"""Branded Google Ads / SEO performance deck generator (python-pptx).

Two layers, mirroring the proven proof-of-concept:

1. RENDER layer (this module's `build_deck`): turns a normalized `DeckData`
   dict into a polished, on-brand .pptx (TBS CI blue #26397A). Pure, deterministic,
   no network — so it is fully unit-testable with sample data.
2. CONTENT layer (`generate_deck_data_from_context`, optional): uses an LLM to
   turn a raw data snapshot (e.g. a Looker Studio PDF digest or SE Ranking data)
   into that normalized dict, applying the client-facing rules (only real data,
   positive-but-honest framing, one topic per slide).

Keep the two separate: the renderer must never invent numbers — it only lays out
what the content layer (or caller) hands it.
"""
from __future__ import annotations

from io import BytesIO
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_LABEL_POSITION
from pptx.oxml.ns import qn


# ---- Brand system (matches the approved demo) ------------------------------
CI = RGBColor(0x26, 0x39, 0x7A)        # primary CI blue
CI_DARK = RGBColor(0x1B, 0x2A, 0x5E)   # deeper shade for dark slides
CI_LIGHT = RGBColor(0xE8, 0xEC, 0xF6)  # pale tint for card fills / table bands
ACCENT = RGBColor(0xC9, 0xA2, 0x27)    # restrained gold accent
INK = RGBColor(0x1A, 0x1F, 0x2B)       # near-black body text
MUTE = RGBColor(0x6B, 0x72, 0x80)      # muted gray captions
LINE = RGBColor(0xD9, 0xDE, 0xE8)      # hairline borders
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xCA, 0xDC, 0xFC)       # light blue on dark slides
SOFT = RGBColor(0x9D, 0xB0, 0xDE)      # muted blue on dark slides

HEAD = "Georgia"   # header font with personality
BODY = "Calibri"   # clean body font

EMU_W, EMU_H = Inches(13.333), Inches(7.5)   # 16:9 widescreen


# ---- Low-level helpers ------------------------------------------------------
def _text(slide, x, y, w, h, text, *, size=14, font=BODY, color=INK, bold=False,
          italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0,
          char_spacing=None):
    """Add a text box. `char_spacing` is letter-spacing in points (tracking)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    first = True
    for line in str(text).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        f = run.font
        f.size = Pt(size)
        f.name = font
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if char_spacing is not None:
            run.font._rPr.set("spc", str(int(char_spacing * 100)))
    return box


def _rect(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp


def _ellipse(slide, x, y, d, *, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _slide(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    rect = _rect(s, -0.05, -0.05, 13.43, 7.6, fill=bg)
    # send background to back
    sp = rect._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def _footer(slide, n, d):
    site = d.get("site_label", "")
    report = d.get("report_label", "Performance Report")
    label = f"{site}  ·  {report}  ·  {d['period']}" if site else f"{report}  ·  {d['period']}"
    _text(slide, 0.5, 7.08, 9, 0.3, label, size=9, color=MUTE)
    if d.get("sample"):
        _text(slide, 10.3, 7.08, 2.5, 0.3, "SAMPLE DATA", size=9, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    _text(slide, 12.6, 7.08, 0.3, 0.3, str(n), size=9, color=MUTE, align=PP_ALIGN.RIGHT)


def _title(slide, kicker, text):
    _text(slide, 0.5, 0.45, 12, 0.3, kicker.upper(), size=12, bold=True, color=CI, char_spacing=2)
    _text(slide, 0.5, 0.74, 12.3, 0.7, text, size=30, font=HEAD, bold=True, color=INK)


def _hbar_chart(slide, x, y, w, h, labels, values, *, value_labels=True):
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("series", values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 60
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CI
    # category axis styling
    cat = chart.category_axis
    cat.tick_labels.font.size = Pt(12)
    cat.tick_labels.font.name = BODY
    cat.format.line.fill.background()
    cat.major_tick_mark = XL_TICK_MARK.NONE
    # value axis hidden
    val = chart.value_axis
    val.visible = False
    val.has_major_gridlines = False
    if value_labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(11)
        dl.font.name = BODY
        dl.font.color.rgb = WHITE
        dl.position = XL_LABEL_POSITION.INSIDE_END
    return chart


def _kpi_table(slide, x, y, w, rows, col_w):
    n = len(rows)
    tbl_shape = slide.shapes.add_table(n, len(rows[0]), Inches(x), Inches(y),
                                       Inches(w), Inches(0.42 * n))
    table = tbl_shape.table
    # disable default banded style
    tbl_shape._element.graphic.graphicData.tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"
    for ci, cw in enumerate(col_w):
        table.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        table.rows[ri].height = Inches(0.42)
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = CI
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else CI_LIGHT
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else (PP_ALIGN.CENTER if len(col_w) > 3 else PP_ALIGN.RIGHT)
            run = p.add_run()
            run.text = str(val)
            f = run.font
            f.size = Pt(12 if ri == 0 else 13)
            f.name = BODY
            f.bold = ri == 0 or ci == 0
            f.color.rgb = WHITE if ri == 0 else INK
    return table


# ---- Slide builders ---------------------------------------------------------
def _cover(prs, d):
    s = _slide(prs, bg=CI)
    _rect(s, 8.6, -0.05, 4.78, 7.6, fill=CI_DARK)
    _text(s, 9.0, 3.1, 4.0, 0.4, d.get("hero_kicker", "EMBROIDERY EXCELLENCE"),
          size=12, color=SOFT, align=PP_ALIGN.CENTER, char_spacing=2)
    _text(s, 8.9, 3.5, 4.2, 1.4, d.get("hero_line", "Precision\nManufacturing"),
          size=24, font=HEAD, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
    _text(s, 0.8, 1.0, 7.5, 0.5, d["company"].upper(), size=16, bold=True, color=WHITE, char_spacing=3)
    _text(s, 0.75, 2.4, 7.6, 2.1, d.get("report_title", "Google Ads\nPerformance Report"),
          size=44, font=HEAD, bold=True, color=WHITE, spacing=1.0)
    _text(s, 0.8, 5.25, 7.5, 0.4, f"Reporting Period:  {d['period']}", size=16, color=ICE)
    _text(s, 0.8, 5.65, 7.5, 0.4, f"Prepared by {d.get('prepared_by', 'TBS Marketing')}", size=13, color=SOFT)
    if d.get("sample"):
        _text(s, 0.8, 6.9, 7.5, 0.3, "DEMO DECK — ALL FIGURES ARE PLACEHOLDER SAMPLE DATA",
              size=10, bold=True, color=ACCENT, char_spacing=1)
    return s


# ---- SEO-shaped slide builders (driven by real SE Ranking data) ------------
def _seo_top_keywords(prs, d):
    s = _slide(prs)
    _title(s, "Keyword Rankings", d.get("keyword_headline", "Where You Rank Today"))
    kws = d["top_keywords"][:8]
    # chart: search volume of top-ranked keywords (real, meaningful magnitude)
    if any(k.get("volume") for k in kws):
        _text(s, 0.5, 1.7, 8.0, 0.35, "Top-Ranked Keywords by Monthly Search Volume",
              size=13, bold=True, color=CI)
        _hbar_chart(s, 0.5, 2.1, 8.0, 4.2,
                    [k["term"] for k in kws], [k.get("volume", 0) for k in kws])
    else:
        rows = [["Keyword", "Position", "Change"]] + [
            [k["term"], k["position"], k["change"]] for k in kws
        ]
        _kpi_table(s, 0.5, 1.7, 8.0, rows, [4.6, 1.7, 1.7])
    hero = d["hero_keyword"]
    _rect(s, 8.8, 1.9, 4.0, 4.1, fill=CI, rounded=True)
    _text(s, 9.1, 2.2, 3.4, 0.3, "BEST-RANKING KEYWORD", size=11, bold=True, color=SOFT, char_spacing=1)
    _text(s, 9.1, 2.6, 3.6, 0.9, hero["term"], size=28, font=HEAD, bold=True, color=WHITE)
    _text(s, 9.1, 3.6, 3.4, 0.9, hero["position"], size=54, font=HEAD, bold=True, color=ACCENT)
    _text(s, 9.1, 4.55, 3.4, 0.3, hero.get("sub", ""), size=13, color=ICE)
    _text(s, 9.1, 5.0, 3.5, 0.8, hero.get("note", ""), size=12, italic=True, color=CI_LIGHT, spacing=1.1)
    _footer(s, 3, d)
    return s


def _seo_movers(prs, d):
    s = _slide(prs)
    _title(s, "Movement This Period", d.get("movers_headline", "Biggest Climbers & Watch List"))
    # Two columns: climbers (left), drops (right)
    climbers = d["climbers"][:6]
    drops = d["drops"][:6]
    _text(s, 0.5, 1.65, 6, 0.35, "▲  Biggest Climbers", size=14, bold=True, color=CI)
    rows_c = [["Keyword", "Now", "▲"]] + [[k["term"], k["position"], k["change"]] for k in climbers]
    _kpi_table(s, 0.5, 2.05, 6.0, rows_c, [3.6, 1.2, 1.2])
    _text(s, 7.0, 1.65, 6, 0.35, "Watch List", size=14, bold=True, color=MUTE)
    rows_d = [["Keyword", "Now", "Δ"]] + [[k["term"], k["position"], k["change"]] for k in drops]
    _kpi_table(s, 7.0, 2.05, 6.0, rows_d, [3.6, 1.2, 1.2])
    _footer(s, 4, d)
    return s


def build_seo_deck(data: Dict) -> bytes:
    """Render an SEO-rankings deck from normalized data (numbers from SE Ranking)."""
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    _cover(prs, data)
    _exec_summary(prs, data, page=2)
    _seo_top_keywords(prs, data)
    _seo_movers(prs, data)
    _insights(prs, data, page=5)
    _closing(prs, data)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _exec_summary(prs, d, page=2):
    s = _slide(prs)
    _title(s, "Executive Summary", d["exec_headline"])
    _text(s, 0.5, 1.55, 12.3, 0.9, d["exec_summary"], size=15, color=INK, spacing=1.15)
    kpis = d["kpis"]
    cols, cw, ch, gx, gy, x0, y0 = 4, 2.92, 1.55, 0.18, 0.22, 0.5, 2.7
    for i, k in enumerate(kpis[:8]):
        cx = x0 + (i % cols) * (cw + gx)
        cy = y0 + (i // cols) * (ch + gy)
        _rect(s, cx, cy, cw, ch, fill=CI_LIGHT, line=LINE, rounded=True)
        _text(s, cx + 0.2, cy + 0.16, cw - 0.4, 0.3, k["label"].upper(), size=11, bold=True, color=CI, char_spacing=1)
        _text(s, cx + 0.2, cy + 0.42, cw - 0.4, 0.6, k["value"], size=28, font=HEAD, bold=True, color=INK)
        _text(s, cx + 0.2, cy + 1.08, cw - 0.4, 0.3, k.get("delta", ""), size=11, color=MUTE)
    _footer(s, page, d)
    return s


def _account(prs, d):
    s = _slide(prs)
    _title(s, "Account Performance", d.get("account_headline", "Where the Budget Worked Hardest"))
    rows = [["Metric", "This Period", "Trend"]] + [[r["metric"], r["value"], r["trend"]] for r in d["account_rows"]]
    _kpi_table(s, 0.5, 1.7, 6.3, rows, [3.0, 1.9, 1.4])
    dev = d["device_split"]
    _text(s, 7.2, 1.7, 5.6, 0.35, "Conversions by Device", size=13, bold=True, color=CI)
    _hbar_chart(s, 7.2, 2.1, 5.6, 3.9, [x["label"] for x in dev], [x["value"] for x in dev])
    _footer(s, 3, d)
    return s


def _keywords(prs, d):
    s = _slide(prs)
    _title(s, "Keyword Conversion Performance", d.get("keyword_headline", "The Searches Driving Real Enquiries"))
    kw = d["keywords"]
    _hbar_chart(s, 0.5, 1.7, 8.0, 4.6, [x["term"] for x in kw], [x["conversions"] for x in kw])
    hero = d["hero_keyword"]
    _rect(s, 8.8, 1.9, 4.0, 4.1, fill=CI, rounded=True)
    _text(s, 9.1, 2.2, 3.4, 0.3, "TOP CONVERTING KEYWORD", size=11, bold=True, color=SOFT, char_spacing=1)
    _text(s, 9.1, 2.6, 3.4, 0.8, hero["term"], size=34, font=HEAD, bold=True, color=WHITE)
    _text(s, 9.1, 3.6, 3.4, 0.9, str(hero["conversions"]), size=54, font=HEAD, bold=True, color=ACCENT)
    _text(s, 9.1, 4.55, 3.4, 0.3, hero.get("sub", ""), size=13, color=ICE)
    _text(s, 9.1, 5.0, 3.5, 0.8, hero.get("note", ""), size=12, italic=True, color=CI_LIGHT, spacing=1.1)
    _footer(s, 4, d)
    return s


def _landing(prs, d):
    s = _slide(prs)
    _title(s, "Landing Page Performance", d.get("landing_headline", "Pages That Convert Visitors Into Enquiries"))
    rows = [["Landing Page", "Users", "Avg. Dwell", "Conversions"]] + [
        [r["page"], r["users"], r["dwell"], r["conversions"]] for r in d["landing_pages"]
    ]
    # table with 4 cols; reuse _kpi_table but it center-aligns when >3 cols
    n = len(rows)
    tbl = slide_table = _kpi_table(s, 0.5, 1.75, 9.0, rows, [4.2, 1.6, 1.6, 1.6])
    sp = d.get("landing_spotlight")
    if sp:
        _rect(s, 9.9, 1.75, 2.9, 3.6, fill=CI_LIGHT, line=LINE, rounded=True)
        _text(s, 10.1, 2.0, 2.5, 0.3, sp["label"].upper(), size=10, bold=True, color=CI)
        _text(s, 10.1, 2.35, 2.5, 0.8, sp["value"], size=44, font=HEAD, bold=True, color=INK)
        _text(s, 10.1, 3.2, 2.5, 0.3, sp.get("sub", ""), size=12, color=MUTE)
        _text(s, 10.1, 3.7, 2.55, 1.4, sp.get("note", ""), size=12, italic=True, color=INK, spacing=1.1)
    _footer(s, 5, d)
    return s


def _demographics(prs, d):
    s = _slide(prs)
    _title(s, "Demographics Performance", d.get("demo_headline", "Who Is Responding Best"))
    g = d["gender_cvr"]
    _text(s, 0.5, 1.7, 6, 0.35, "Conversion Rate by Gender", size=13, bold=True, color=CI)
    _hbar_chart(s, 0.5, 2.1, 6.2, 4.0, [x["label"] for x in g], [x["value"] for x in g])
    _text(s, 7.2, 1.7, 5.6, 0.35, "Top Age Segments", size=13, bold=True, color=CI)
    for i, a in enumerate(d["age_segments"][:4]):
        cy = 2.15 + i * 0.95
        _rect(s, 7.2, cy, 5.6, 0.78, fill=CI_LIGHT, line=LINE, rounded=True)
        _text(s, 7.45, cy, 2.5, 0.78, a["label"], size=20, font=HEAD, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 9.8, cy, 2.8, 0.78, a["value"], size=22, font=HEAD, bold=True, color=CI,
              anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    _footer(s, 6, d)
    return s


def _insights(prs, d, page=7):
    s = _slide(prs)
    _title(s, "Strategic Insights", d.get("insights_headline", "Where to Focus Next Month"))
    recs = d["recommendations"]
    cw, ch, gx, gy, x0, y0 = 6.0, 2.1, 0.3, 0.3, 0.5, 1.75
    for i, r in enumerate(recs[:4]):
        cx = x0 + (i % 2) * (cw + gx)
        cy = y0 + (i // 2) * (ch + gy)
        _rect(s, cx, cy, cw, ch, fill=WHITE, line=CI, line_w=1.25, rounded=True)
        _text(s, cx + 0.25, cy + 0.2, 1.2, 0.9, f"{i+1:02d}", size=40, font=HEAD, bold=True, color=CI_LIGHT)
        _text(s, cx + 1.3, cy + 0.28, cw - 1.5, 0.5, r["title"], size=19, font=HEAD, bold=True, color=INK)
        _text(s, cx + 1.3, cy + 0.85, cw - 1.55, 1.1, r["body"], size=13, color=INK, spacing=1.12)
    _footer(s, page, d)
    return s


def _closing(prs, d):
    s = _slide(prs, bg=CI)
    _rect(s, -0.05, -0.05, 13.43, 0.12, fill=ACCENT)
    _text(s, 0.8, 1.2, 8, 0.4, "KEY TAKEAWAYS", size=13, bold=True, color=SOFT, char_spacing=3)
    _text(s, 0.75, 1.7, 9, 1.6, d.get("closing_headline", "Positive Momentum,\nClear Next Moves"),
          size=40, font=HEAD, bold=True, color=WHITE)
    for i, p in enumerate(d["takeaways"][:3]):
        cy = 3.7 + i * 0.6
        _ellipse(s, 0.85, cy + 0.06, 0.16, fill=ACCENT)
        _text(s, 1.2, cy, 9.5, 0.5, p, size=16, color=CI_LIGHT)
    _text(s, 0.8, 5.9, 6, 0.7, "Thank You", size=30, font=HEAD, italic=True, bold=True, color=ACCENT)
    _text(s, 0.8, 6.6, 9, 0.4, f"{d.get('prepared_by', 'TBS Marketing')}  ·  Prepared for {d['company']}",
          size=13, color=SOFT)
    return s


# ---- Public API -------------------------------------------------------------
def build_deck(data: Dict) -> bytes:
    """Render the normalized deck data dict into .pptx bytes."""
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H

    _cover(prs, data)
    _exec_summary(prs, data)
    _account(prs, data)
    _keywords(prs, data)
    _landing(prs, data)
    _demographics(prs, data)
    _insights(prs, data)
    _closing(prs, data)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def sample_deck_data() -> Dict:
    """Placeholder data mirroring the approved demo — for tests and previews.

    Every value here is fictional sample data, clearly flagged on the cover.
    """
    return {
        "sample": True,
        "company": "MicroEmbroidery",
        "site_label": "MicroEmbroidery.com",
        "report_label": "Google Ads Performance",
        "report_title": "Google Ads\nPerformance Report",
        "period": "1 – 22 May 2026",
        "prepared_by": "TBS Marketing",
        "hero_kicker": "งานปัก  ·  EMBROIDERY EXCELLENCE",
        "hero_line": "Precision\nManufacturing",
        "exec_headline": "A Strong Month of Measurable Growth",
        "exec_summary": (
            "Paid search delivered efficient, conversion-focused growth in early May. "
            "Conversions rose while cost per conversion fell, and click-through rate "
            "strengthened — signalling sharper targeting and healthy demand for embroidery "
            "and patch manufacturing."
        ),
        "kpis": [
            {"label": "Spend", "value": "฿48,250", "delta": "Efficient pacing"},
            {"label": "Conversions", "value": "63", "delta": "▲ 18% vs. prior"},
            {"label": "Cost / Conv.", "value": "฿766", "delta": "▼ 12% — more efficient"},
            {"label": "CTR", "value": "6.4%", "delta": "▲ 0.9 pts"},
            {"label": "Clicks", "value": "2,940", "delta": "▲ 11%"},
            {"label": "Impressions", "value": "45,900", "delta": "Steady reach"},
            {"label": "Avg. CPC", "value": "฿16.4", "delta": "Stable"},
            {"label": "Conv. Rate", "value": "2.1%", "delta": "▲ 0.3 pts"},
        ],
        "account_rows": [
            {"metric": "Spend", "value": "฿48,250", "trend": "On plan"},
            {"metric": "Clicks", "value": "2,940", "trend": "▲ 11%"},
            {"metric": "Impressions", "value": "45,900", "trend": "▲ 4%"},
            {"metric": "CTR", "value": "6.4%", "trend": "▲ 0.9 pts"},
            {"metric": "Avg. CPC", "value": "฿16.4", "trend": "Stable"},
            {"metric": "Conversions", "value": "63", "trend": "▲ 18%"},
            {"metric": "Cost / Conversion", "value": "฿766", "trend": "▼ 12%"},
        ],
        "device_split": [
            {"label": "Mobile", "value": 38},
            {"label": "Desktop", "value": 21},
            {"label": "Tablet", "value": 4},
        ],
        "keywords": [
            {"term": "งานปัก", "conversions": 17},
            {"term": "embroidered patch", "conversions": 12},
            {"term": "custom badge", "conversions": 9},
            {"term": "iron-on patch", "conversions": 7},
            {"term": "logo embroidery", "conversions": 6},
            {"term": "applique maker", "conversions": 4},
        ],
        "hero_keyword": {
            "term": "งานปัก", "conversions": 17, "sub": "conversions  ·  9.1% CTR",
            "note": "Highest-intent term this period — strong candidate to scale.",
        },
        "landing_pages": [
            {"page": "/ (Homepage)", "users": "501", "dwell": "1:42", "conversions": "24"},
            {"page": "/custom-embroidery-patches", "users": "212", "dwell": "2:18", "conversions": "18"},
            {"page": "/badges-appliques", "users": "146", "dwell": "1:55", "conversions": "11"},
            {"page": "/contact-quote", "users": "98", "dwell": "2:40", "conversions": "7"},
            {"page": "/export-manufacturing", "users": "73", "dwell": "1:30", "conversions": "3"},
        ],
        "landing_spotlight": {
            "label": "Homepage Spotlight", "value": "501", "sub": "users this period",
            "note": "Strongest entry point and top converter — a reliable foundation to build campaigns around.",
        },
        "gender_cvr": [
            {"label": "Undetermined", "value": 16.6},
            {"label": "Female", "value": 8.2},
            {"label": "Male", "value": 5.9},
        ],
        "age_segments": [
            {"label": "25–34", "value": "31%"},
            {"label": "35–44", "value": "27%"},
            {"label": "45–54", "value": "18%"},
            {"label": "18–24", "value": "14%"},
        ],
        "recommendations": [
            {"title": "Scale the winners", "body": "Increase budget on “งานปัก” and “embroidered patch” — both convert above account average at efficient cost."},
            {"title": "Tighten keyword targeting", "body": "Add high-intent modifiers (custom, bulk, export) and trim broad terms with clicks but no conversions."},
            {"title": "Reallocate to mobile", "body": "Mobile drives 60% of conversions — shift bid adjustments and budget weighting accordingly."},
            {"title": "Refine audience bidding", "body": "Layer the 25–44 age segment as a bid boost where it consistently outperforms."},
        ],
        "takeaways": [
            "Conversions up 18% at a 12% lower cost per conversion",
            "“งานปัก” confirmed as the top-intent, scalable keyword",
            "Mobile and the 25–44 audience are the growth engines",
        ],
        "closing_headline": "Positive Momentum,\nClear Next Moves",
    }


def sample_seo_deck_data() -> Dict:
    """Placeholder SEO data shaped like the data→deck bridge output — for tests."""
    return {
        "sample": True,
        "company": "Cleanatic",
        "site_label": "cleanatic.com",
        "report_label": "SEO Performance",
        "report_title": "SEO Performance\nReport",
        "period": "May 2026",
        "prepared_by": "TBS Marketing",
        "hero_kicker": "ORGANIC SEARCH",
        "hero_line": "Visibility\n& Growth",
        "exec_headline": "Steady Gains in Organic Visibility",
        "exec_summary": (
            "Organic rankings strengthened across the tracked keyword set this month. "
            "More terms moved onto page one and average position improved, reflecting "
            "healthy momentum from ongoing optimisation work."
        ),
        "kpis": [
            {"label": "Keywords Tracked", "value": "48", "delta": "Full set"},
            {"label": "On Page 1", "value": "21", "delta": "▲ pos 1–10"},
            {"label": "Avg. Position", "value": "7.5", "delta": "▲ improving"},
            {"label": "Improved", "value": "19", "delta": "this period"},
        ],
        "top_keywords": [
            {"term": "บริการทำความสะอาด", "position": "3", "change": "▲2", "volume": 5400},
            {"term": "แม่บ้านรายวัน", "position": "5", "change": "▲1", "volume": 3300},
            {"term": "ทำความสะอาดบ้าน", "position": "6", "change": "▲4", "volume": 2900},
            {"term": "cleaning service bangkok", "position": "8", "change": "▲3", "volume": 1600},
            {"term": "แม่บ้านมืออาชีพ", "position": "9", "change": "—", "volume": 1200},
            {"term": "deep cleaning", "position": "11", "change": "▼1", "volume": 880},
        ],
        "hero_keyword": {
            "term": "บริการทำความสะอาด", "position": "#3", "sub": "5,400 searches / mo",
            "note": "Your strongest high-volume term — a click away from the top spots.",
        },
        "climbers": [
            {"term": "ทำความสะอาดบ้าน", "position": "6", "change": "▲4"},
            {"term": "cleaning service bangkok", "position": "8", "change": "▲3"},
            {"term": "บริการทำความสะอาด", "position": "3", "change": "▲2"},
            {"term": "แม่บ้านรายวัน", "position": "5", "change": "▲1"},
        ],
        "drops": [
            {"term": "deep cleaning", "position": "11", "change": "▼1"},
            {"term": "office cleaning", "position": "18", "change": "▼3"},
        ],
        "recommendations": [
            {"title": "Push page-1 near-misses", "body": "“deep cleaning” sits at #11 — targeted internal links and content refresh can move it onto page one."},
            {"title": "Defend the winners", "body": "Keep content fresh for the high-volume terms now in the top 5 to hold and improve position."},
            {"title": "Expand winning topics", "body": "Build supporting content around the climbing cleaning-service cluster to capture related queries."},
            {"title": "Review the watch list", "body": "Investigate the small declines on office-cleaning terms before they slip further."},
        ],
        "takeaways": [
            "21 of 48 keywords now rank on page one",
            "Average position improved to 7.5",
            "High-volume “บริการทำความสะอาด” climbed to #3",
        ],
        "closing_headline": "Building Momentum\nin Organic Search",
    }


if __name__ == "__main__":
    out = build_deck(sample_deck_data())
    with open("GoogleAds_Report_PY.pptx", "wb") as fh:
        fh.write(out)
    print(f"Wrote GoogleAds_Report_PY.pptx ({len(out)} bytes)")

    seo = build_seo_deck(sample_seo_deck_data())
    with open("SEO_Report_PY.pptx", "wb") as fh:
        fh.write(seo)
    print(f"Wrote SEO_Report_PY.pptx ({len(seo)} bytes)")
