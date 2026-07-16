"""AI-designed deck generator — free engine, downloadable file.

Pipeline (no third-party deck subscription):
  1. assemble real data (done elsewhere) -> a data brief
  2. an editable, Abacus-style PROMPT + the data -> an LLM (DeepSeek/Groq) writes a
     complete, self-contained HTML presentation (unique design each time)
  3. headless Chromium renders that HTML to the chosen downloadable format:
       - PDF  (pixel-perfect, one slide per page)
       - PPTX (each slide screenshotted full-bleed into a python-pptx slide)

The AI designs the slides (not a fixed template); we only render its HTML to a file.
"""
from __future__ import annotations

import logging
import asyncio
import base64
import re
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Optional

from services.ai_service import ai_service, ProgressCb

logger = logging.getLogger(__name__)

# Slide geometry — 16:9 at 1920x1080 CSS px, mapped to a 13.333x7.5in PPTX/PDF page.
SLIDE_W_PX, SLIDE_H_PX = 1920, 1080
# Render at this device-pixel-ratio so the screenshots are high-res (sharp in the
# PDF/PPTX). The deck is AUTHORED in 1920x1080 CSS px — only the raster scales
# (2x -> 2560x1440). PDF DPI scales with it so the physical page size stays 13.33x7.5in.
_RENDER_SCALE = 1.5

# AI photos per deck. Bounds both the early (streamed) prewarm and the final resolve,
# and the concurrency cap keeps us under Supermachine's 10 req/min rate limit.
MAX_DECK_IMAGES = 8
_IMG_CONCURRENCY = 3


def _img_prompt(tag: str) -> str:
    """Extract the data-prompt text from an <img class="ai-img" ...> placeholder."""
    m = re.search(r'data-prompt=["\'](.*?)["\']', tag, re.S)
    return (m.group(1).strip() if m else "professional abstract background, premium")


# ---------------------------------------------------------------------------
# The editable prompt. This is the default; callers may pass their own text.
# {brand}, {structure}, {data} are filled in at request time.
# The HTML OUTPUT CONTRACT at the end is required for rendering and should be
# kept even if the rest of the prompt is customised.
# ---------------------------------------------------------------------------
DEFAULT_DECK_PROMPT = """You are an elite, executive-level Data Scientist, Presentation Designer, and Growth Strategist. Your task is to transform the attached monthly performance data into a visually stunning, dynamic, and premium presentation deck.

The report must look like it was hand-crafted by a premium data-analytics agency for a C-suite executive review. It must be clean, highly modern, clear, and deeply strategic.

=======================================================
CRITICAL RULES FOR DATA ACCURACY & TONE:
1. Use ONLY the actual data provided inside the file/text. Do NOT add fake numbers, assumptions, or estimates.
2. REPORT DECLINES HONESTLY AND PROMINENTLY. If a metric, query, page or theme dropped, say so plainly — state the actual number and movement (e.g. "clicks 1,240 -> 980, -21%"), the likely cause, and the specific fix. Never hide, omit, bury or spin a decline into vague positivity. Keep the tone professional and calm (a senior consultant briefing a client), NOT alarmist — but the client must be able to see exactly what is going down and what you will do about it. Anything flagged in ANALYST FLAGS as AT RISK / DEFEND / PAGE DECLINE MUST appear in the deck.
3. Follow the structural rule: One core narrative or trend = one slide. Do NOT overcrowd layout spaces with dense blocks of text.

=======================================================
DYNAMIC BRANDING & STYLE ENGINE — EDITORIAL / DESIGN-STUDIO GRADE:
Design this like a high-end editorial design studio "frame system" — think art-direction-led poster design and a premium print magazine, NOT a generic corporate slide template. Every deck must feel like a confident, hand-crafted design object with its own identity. Analyze the domain and data context to establish that identity:
- INDUSTRY VIBE: Deduce the company's core focus (e.g., Trust/Security for Insurance; Precision for Manufacturing) and let it drive a DISTINCTIVE art direction — two decks must look visibly different.
- EXPRESSIVE DISPLAY TYPOGRAPHY (the #1 driver of the look): Pair a characterful display face (an elegant high-contrast SERIF, or a bold grotesque) for headlines with a clean sans for body — load real fonts from fonts.googleapis.com (e.g. Fraunces, Playfair Display, Libre Caslon, Archivo, Space Grotesk, Instrument Serif, Bricolage Grotesque, Syne). Headlines are OVERSIZED and confident (clamp ~84–156px), tight line-height, with deliberate mixing of an italic serif accent word inside a sans headline (e.g. big sans + one *italic serif* word). Use real typographic hierarchy: tiny ALL-CAPS letter-spaced kickers/eyebrows, large display headline, restrained body.
- KICKER / SYSTEM LABEL MOTIF: Put a small uppercase, letter-spaced label in a corner of each slide (e.g. "PERFORMANCE SYSTEM · VOL. 01", "ORGANIC SEARCH — 2026", or a slide index "03 / 09"). This "frame system" tagging is a signature of the look.
- THE 60-30-10 COLOR RULE:
  * 60% Dominant: a premium CREAM / off-white / ivory paper ground (not stark white) for editorial warmth and readability — or, for cover/section/closing slides, ONE saturated full-bleed colour field.
  * 30% Secondary Structure: deep near-black ink, charcoal, or a deep brand tone for crisp text hierarchy.
  * 10% Accent/Pop Color: exactly ONE vivid, sophisticated accent (a confident yellow, vermilion red, electric blue, etc.) used for emphasis, the kicker, key metrics and a single decorative motif.
- DECORATIVE MOTIF (use with RESTRAINT, never over charts/tables): carry ONE repeating graphic device across slides — e.g. a thin ruled grid or dot-grid, a slightly rotated/"crooked" outlined frame, a small filled star or circle, pill/tape-style labels, a hairline rule with an index number. It should feel intentional and editorial, never like clip-art.
- LAYOUT: confident, asymmetric, magazine-grade composition with generous negative space and a clear focal point — avoid centred, evenly-stacked "default template" layouts. Vary layouts slide to slide (full-bleed type poster cover, split editorial spreads, a strong data spread).
- COVER, SECTION DIVIDERS & CLOSING: treat these as poster pages — huge expressive type on a bold colour field or dramatic photo, minimal supporting text. These set the tone.

=======================================================
INTERACTIVE DATA VISUALIZATION (PLOTLY.JS INTEGRATION):
Wherever the data can be grouped, tracked over time, or categorized into distributions, add a chart.
- Charts are embedded ONLY through the hidden chart mechanism defined in the HTML OUTPUT CONTRACT
  below (a sized <div> + an immediately-following <script type="application/json" class="plotly-spec">).
  NEVER write the chart's JSON/config anywhere a reader can see it — no visible JSON, no <pre>/<code>
  block, no "config" column, no spec printed next to the chart. The JSON lives only inside that script.
- Charts strictly inherit the chosen palette (secondary tone for gridlines/text, the accent for the
  key series). Keep them clean: few gridlines, clean tooltips, sharp modern fonts.

=======================================================
DYNAMIC DECK STRUCTURE:
Do not limit the output to a rigid slide count. Read the data payload and dynamically generate a
beautifully balanced sequence of slides mapping to the structure below.

CRITICAL OUTPUT RULE: emit each slide as FINAL RENDERED HTML only. Do NOT print any planning text,
slide blueprints, "Slide X:" headings, "VISUAL LAYOUT COMPOSITION", or chart-config blocks as visible
content — translate all of that directly into the HTML/CSS and the hidden chart script. The reader must
see a finished slide, never the instructions or data behind it.

{structure_directive}

=======================================================
INPUT DATA TO PROCESS:
DATA (use ONLY this — pulled automatically from the connected data source for this site):
{data}"""


# Always appended by build_prompt — the renderer depends on it, so it stays out of
# the user-editable prompt and can't be removed by mistake.
HTML_CONTRACT = """=== HTML OUTPUT CONTRACT (required — this OVERRIDES any conflicting instruction above) ===
Output ONE complete, self-contained HTML document and NOTHING ELSE: no markdown, no commentary, and do NOT print slide specifications, "SLIDE NUMBER & TITLE", "VISUAL LAYOUT COMPOSITION", or standalone code blocks as text — translate ALL of that directly into the final rendered HTML.
- NEVER show raw chart JSON or any {"data":...,"layout":...} config as visible text. A chart's JSON belongs ONLY inside its hidden <script type="application/json" class="plotly-spec"> — if a reader can see JSON on a slide, that is a bug. Each chart spec must appear EXACTLY ONCE, inside that one hidden script; do NOT also place a copy in a visible <pre>/<code>/<div> or as a text node, and do NOT HTML-escape a copy onto the slide.
- KPI DISCIPLINE — do NOT repeat the same headline metric row on every slide. Show the clicks / impressions / CTR / avg-position (or sessions/cost) KPI strip on the EXECUTIVE SUMMARY slide ONLY; every other slide focuses on its own chart/insight with at most one or two metrics relevant to THAT slide. The cover and closing carry NO KPI chips.
- EVERY slide must fill the entire 1920x1080 page edge to edge (a real full-bleed page): no letterboxing, no large empty margins or blank bands — size the hero chart/visual to consume the slide's main area.
- Start with <!DOCTYPE html>.
- All CSS inline in a single <style> tag. Fonts may load from fonts.googleapis.com.
- Each slide is exactly: <section class="slide"> ... </section>. Every .slide is 1920px wide and 1080px tall, overflow hidden.
- FILL THE CANVAS — this is critical. Every slide MUST use the full 1920x1080 height; NO slide may leave a large empty band (no more than ~12% blank vertical space). Make each .slide a vertical flex container (display:flex; flex-direction:column; justify-content:center; gap:40px; padding:88px 108px) so content is balanced over the whole height — never dump everything in the top third. To fill space meaningfully: pair every data TABLE or KPI-card row with a relevant Plotly chart beside or below it (two-column or stacked grid), enlarge cards, and use generous spacing. A data slide that is only a small table at the top is NOT acceptable — add a chart or distribute the layout to fill the slide.
- Include this CSS so it paginates when printed:
    @page { size: 1920px 1080px; margin: 0; }
    * { box-sizing: border-box; }
    html,body { margin:0; padding:0; }
    .slide { width:1920px; height:1080px; position:relative; overflow:hidden; page-break-after:always; }
- CHARTS via Plotly.js (use them generously where the data suits a chart). The Plotly library is provided for you — do NOT add any Plotly <script src>, and do NOT call Plotly.newPlot yourself. For each chart output exactly two elements:
    (a) a sized container in its slide: <div id="chartN" style="width:900px;height:570px"></div>  (unique id per chart)
    (b) immediately AFTER it, its spec: <script type="application/json" class="plotly-spec" data-target="chartN">{"data":[...],"layout":{...}}</script>
  The system reads every plotly-spec and renders it. The JSON must be strictly valid (double-quoted keys/strings, no trailing commas, no JS expressions). Never print chart config as visible text. Charts MUST fit fully within the 1920x1080 slide.
  Every chart's "layout" MUST be styled consistently for a premium look: "paper_bgcolor":"rgba(0,0,0,0)","plot_bgcolor":"rgba(0,0,0,0)", "font":{"family":"<your body font>","color":"<--muted>","size":13}, "margin":{"t":10,"l":48,"r":16,"b":36}, "showlegend":false (or a single clean legend), x/y axes with "gridcolor":"<--line>","zeroline":false and no chartjunk. Use the ACCENT colour for the primary series and --accent-2 for a secondary series — never default Plotly colours. Keep it clean: few gridlines, sharp modern fonts, rounded bar feel.
- PHOTOS: to add a photo, output <img class="ai-img" data-prompt="<a vivid, specific photographic description on-theme to this site's industry — e.g. 'close-up of embroidered military patches on dark fabric, studio lighting, premium'>" style="...">. Do NOT put a src — the system fills it in. Size/position it with CSS (object-fit:cover) so it fits inside its slide.
  Use photos on MOST slides as full-bleed backgrounds or side panels to make the deck visually rich. When text or numbers sit on top of a photo, place a dark gradient overlay (e.g. linear-gradient(rgba(15,23,42,.65), rgba(15,23,42,.35))) over the image for legibility. Keep chart and table areas on a clean solid background, not over a photo. Use at most 8 photos total.
- ICONS: add a crisp icon with <i class="ai-icon" data-icon="NAME"></i> (no other attributes). The system swaps it for an inline SVG that inherits the current text colour (set it to --accent or --ink via CSS; size with font-size/width). Put ONE relevant icon on every KPI, every list bullet, and section markers — never as random decoration. Available NAMEs ONLY: trending-up, trending-down, target, search, eye, mouse-pointer-click, users, globe, bar-chart, line-chart, pie-chart, activity, arrow-up-right, arrow-down-right, check-circle, alert-triangle, lightbulb, rocket, star, award, flag, calendar, clock, map-pin, link, file-text, layers, filter, zap, dollar-sign, percent, shopping-cart, smartphone, monitor, thumbs-up, refresh-cw, compass, megaphone, sparkles, gauge. If none fit, omit the icon.
- Do NOT reference any other external resource or remote <img> with a real URL. Self-contained except the Google Font, the single Plotly script, and ai-img / ai-icon placeholders the system fills."""


DEFAULT_BRAND = (
    "Primary CI color: #26397A. Style: premium B2B aesthetic — trustworthy, professional. "
    "White or light backgrounds with strong spacing and clean layouts. Modern typography, "
    "consistent hierarchy, subtle CI-blue accents. Avoid overly colorful or playful designs. "
    "Keep visuals sharp, structured, and minimal."
)

# A brand directive that makes the AI invent a UNIQUE visual identity per site,
# rather than reusing one fixed template.
UNIQUE_STYLE_BRAND = (
    "Design like a high-end EDITORIAL DESIGN STUDIO — art-direction-led, poster/magazine grade — NOT a generic "
    "corporate slide template. Create a UNIQUE visual identity tailored specifically to THIS website and its industry; "
    "do NOT reuse a generic or default template, and do not default to navy/blue. From the site's domain and content, choose: "
    "a distinctive colour palette (one dominant colour ~60% — favour a warm cream/ivory paper ground for content slides — "
    "1-2 supporting tones, ONE sharp vivid accent); an EXPRESSIVE display font pairing (a characterful serif or bold grotesque "
    "for OVERSIZED confident headlines + a clean sans for body, loaded from fonts.googleapis.com), mixing an italic serif accent "
    "word into sans headlines; tiny ALL-CAPS letter-spaced kicker/eyebrow labels and a small corner 'system' tag or slide index; "
    "and ONE repeating editorial motif used with restraint (a dot/ruled grid, a slightly rotated outlined frame, a star/circle, "
    "pill/tape labels) carried across every slide — never over charts/tables. Layouts must be confident and ASYMMETRIC with "
    "generous negative space and a clear focal point; vary them slide to slide. Two different sites must produce visibly "
    "different decks. Treat the cover, section dividers and closing as poster pages: huge expressive type on a bold saturated "
    "or dark colour field. Avoid AI-slop: no thin accent lines under titles, no full-width decorative colour bars unless they "
    "serve the layout, no centred evenly-stacked default layouts."
)

DEFAULT_STRUCTURE = (
    "Cover (company, report title, reporting period, hero visual); Executive Summary (CHART-LED — "
    "one or two large charts as the focus, KPIs only as a slim compact strip, not large number cards "
    "with empty space; strongest positives + short strategic summary); Keyword Rankings overview; Biggest Movers; "
    "Wins of the Month; Opportunities; Strategic Recommendations; Closing (key takeaways + thank you)."
)


# Defaults tuned for a Google Ads report built from an uploaded Looker Studio PDF.
GOOGLE_ADS_BRAND = (
    "Primary CI color: #26397A. Style: premium B2B manufacturing / export company "
    "aesthetic (embroidery badges, patches, appliqués). Feel: trustworthy, industrial "
    "precision, professional, global export standard. White or light backgrounds with "
    "strong spacing and clean layouts. Modern typography, consistent hierarchy, subtle "
    "CI-blue accents. Avoid overly colorful or playful designs. Keep visuals sharp, "
    "structured, and minimal."
)

GOOGLE_ADS_STRUCTURE = (
    "1. Cover Slide — company name, 'Google Ads Performance Report', and the REPORTING PERIOD date range as the subtitle, on a professional hero visual. NO KPI numbers/metric chips/clicks/spend on the cover — keep it clean.\n"
    "2. Executive Summary — CHART-LED slide: one or two large Plotly charts as the focus, with the high-level KPIs as a slim compact strip along the top or bottom (NOT large number cards with empty space). Strongest positive outcomes, short strategic summary.\n"
    "3. Account Performance — Spend, Clicks, Impressions, CTR, CPC, Conversions, Cost per Conversion, ROAS/Conversion Value if available; highlight key wins visually.\n"
    "4. Keyword Conversion Performance — top converting keywords, best CTR keywords, highest-value search intent; tables or charts.\n"
    "5. Landing Page Performance — best-performing and conversion-driving pages, engagement insights if available; focus on strengths.\n"
    "6. Conversion Funnel Analysis — funnel visualization, drop-off points, positive conversion flow insights.\n"
    "7. Click Type Analysis — breakdown of click types, performance comparison, strongest engagement sources.\n"
    "8. Demographics Performance — age, gender, audience segments if available; highlight highest-performing demographics.\n"
    "9. Needs Attention — What's Dropping (REQUIRED — always include this slide) — every declining metric, campaign, keyword or landing page RANKED BY IMPACT, each with its real movement (previous → current, % change), the likely cause, and the specific fix (bids, budget, targeting, creative). Calm and professional, never alarmist, but never softened into vague positivity. OMIT only if genuinely nothing declined.\n"
    "10. Strategic Insights & Optimization Opportunities — actionable Google Ads recommendations ONLY (scale winning campaigns, improve keyword targeting, budget allocation, audience optimization, bid strategy). Do NOT suggest website redesign unless the data supports it.\n"
    "11. Closing Slide — key takeaways, honest momentum summary, professional thank-you page with the reporting period in a slim footer."
)


# Structure for the combined monthly report (Google Search Console + Google Analytics).
# When a site has no matching GA4 property the brief carries no analytics data, so the
# GA4-only slides are explicitly OMITTABLE and the deck gracefully degrades to SEO-only.
GSC_STRUCTURE = (
    "1. Cover Slide — site/domain, 'Search Console Report', and the REPORTING PERIOD date range as the subtitle, on a professional hero visual. NO KPI numbers/metric chips/clicks/impressions on the cover — keep it clean.\n"
    "2. Executive Summary — CHART-LED slide: one or two large Plotly charts (clicks/impressions/CTR/position) as the focus, with the high-level KPIs as a slim compact strip along the top or bottom (NOT large number cards with empty space). Strongest positives, short strategic summary.\n"
    "3. Search Performance — clicks, impressions, CTR, average position with period-over-period change; highlight key wins visually.\n"
    "4. Performance Over Time — use MONTHLY PERFORMANCE: a Plotly COMBO chart with clicks & impressions as bars (accent / accent-2) and average POSITION as a line on a SECONDARY y-axis that is REVERSED (lower is better, so a rising line = improving rank). This chart MUST be self-explanatory: include a VISIBLE LEGEND that names each series (e.g. \"Clicks\", \"Impressions\", \"Avg position\") with \"showlegend\":true positioned at the top (and enough top margin that it is NOT clipped), give the primary y-axis the title \"clicks / impressions\" and the secondary y-axis the title \"avg position (lower is better)\", so a reader instantly knows the bars are traffic and the line is ranking. Optionally a second slide with the daily impressions and URL-clicks as filled AREA charts.\n"
    "5. Top Queries (REQUIRED — always include this slide) — a clean, readable TABLE of the top ~10 queries by clicks "
    "with columns Query · Clicks · Impressions · CTR · Avg position. The boss must be able to read exactly which keywords drive "
    "traffic, so this is a real labelled table (NOT just the bubble chart). Optionally pair it with a small bar of clicks by query.\n"
    "6. Keyword Landscape (REQUIRED — always include this slide) (mix summary — NO bubble chart on this slide) — a .layout-split full-height slide. "
    "ONE side: the eyebrow 'KEYWORD LANDSCAPE' + a title, the TOP 3 queries (by clicks) as compact cards "
    "(query · clicks · avg position), and a grounded 'Unique Keywords' .kpi from KEYWORD MIX showing the "
    "'Unique keywords (distinct queries)' number. OTHER side: a ranking-distribution DONUT (Plotly pie, \"hole\":0.6, "
    "values = [positions 1-3 count, positions 4-10 count, positions 11+ count] from KEYWORD MIX, "
    "labels = [\"Pos 1-3\",\"Pos 4-10\",\"Pos 11+\"], colours = [accent, accent-2, --muted], sized to fill its side with a slim "
    "legend below — it MUST actually render as a pie and must NOT overflow/clip. Do NOT replace the donut with plain number "
    "cards, and do NOT print any '% long-tail/short-tail' framing.\n"
    "7. Keyword Opportunity · Position vs Impressions (REQUIRED — always include this slide) (FULL-PAGE bubble chart — this whole slide is the chart) — a slim eyebrow "
    "'KEYWORD OPPORTUNITY' + a short title at the top, then ONE big bubble chart div that fills the REST of the slide (full content "
    "width ~1150px and ~540px tall) as the hero — no side panels, no cards, nothing else competing for space. The chart is ONE "
    "Plotly scatter trace with \"type\":\"scatter\", x = the average position of every query (x-axis REVERSED so better ranks sit on "
    "the RIGHT), y = the impressions of every query. The trace MUST carry, ALIGNED to x/y, a \"customdata\" array = the QUERY STRING "
    "of each point, and MUST include \"meta\":{\"chart\":\"keyword-bubble\"} so the renderer can finish it. Marker colour = --accent. "
    "(The system deterministically scales the bubbles by impressions and labels the biggest ~8 queries beside their bubble, so just "
    "supply accurate x, y and customdata — do not worry about marker.size yourself.) \"showlegend\":false.\n"
    "8. Biggest Movers (REQUIRED: show BOTH risers AND fallers — not climbers only) — two side-by-side panels: rising vs falling "
    "QUERIES (previous → current) and, below or beside them, rising vs falling LANDING PAGES. Frame the fallers honestly as "
    "at-risk/decline-to-defend (not alarmist), but they MUST be shown so the report is balanced. Diverging horizontal bars or clean "
    "two-column tables work well.\n"
    "9. Query Opportunities — high-impression / low-CTR or near-page-1 queries to target next.\n"
    "10. CTR Opportunities — from CTR OPPORTUNITIES: queries with strong impressions but a CTR below what their rank should earn. Show actual vs expected CTR (a grouped/diverging bar or a clean table) ranked by missed clicks, framed as quick title/meta wins. OMIT this slide entirely if CTR OPPORTUNITIES is (none).\n"
    "11. Query Footprint — from QUERY FOOTPRINT: a Plotly STACKED BAR per month of top-10 query counts (pos 1-3 stacked with pos 4-10) plus a line for total queries; shows ranking visibility growing over time.\n"
    "12. Top Pages (REQUIRED — always include this slide) — a readable TABLE of the best-performing landing pages by clicks "
    "(Page · Clicks · Impressions · CTR · Avg position); shorten long URLs to their path. Focus on strengths.\n"
    "13. Channels — Devices & Search Type — a .layout-split: on one side a device breakdown (from BY DEVICE — donut or bar of clicks/impressions by desktop/mobile/tablet); on the other a search-surface breakdown (from BY SEARCH TYPE — bar of web/image/video/news). If BY SEARCH TYPE is (none) or web-only, show devices alone full-width. OMIT the slide only if BY DEVICE is also (none).\n"
    "14. Search Appearance — from SEARCH APPEARANCE: a bar or table of rich-result types (FAQ, product snippets, etc.) by clicks/impressions/CTR. OMIT this slide entirely if SEARCH APPEARANCE is (none).\n"
    "15. Website Audience & Engagement (OMIT this entire slide if there is no WEBSITE ANALYTICS / GA4 section in the data) — from AUDIENCE & ENGAGEMENT: a CHART-LED analytics slide using Google Analytics on-site behaviour. Show sessions, total users, new users, engagement rate, bounce rate, avg session duration and conversions with their period-over-period change as a slim KPI strip, and make the SESSIONS OVER TIME daily series the hero chart (an area/line of sessions & users, with conversions if present). This is the GA4 counterpart to the search-performance slide — clearly label it as website analytics so it's not confused with Search Console clicks.\n"
    "16. Traffic by Channel (OMIT this entire slide if there is no WEBSITE ANALYTICS / GA4 section) — from TRAFFIC BY CHANNEL: how sessions/users/conversions split across channels (organic, direct, paid, referral, social, etc.). Use a horizontal bar or donut of sessions by channel paired with a small table; highlight the strongest acquisition channels and frame conversions positively.\n"
    "17. Geographic Distribution — from GEOGRAPHY: a Plotly CHOROPLETH world map (\"type\":\"choropleth\") shaded by the stated metric (sessions or clicks) using the stated locationmode, paired with a top-countries bar. Keep the geo clean (accent colourscale, transparent bg, no coastline/frame chartjunk).\n"
    "18. Needs Attention — What's Dropping (REQUIRED — always include this slide) — the honest counterpart to the wins. Every declining query, landing page and theme, RANKED BY IMPACT, each with its real movement (previous → current, and the % or position change), the likely cause, and the specific fix. Source this from the ANALYST FLAGS in the data (AT RISK / DEFEND / PAGE DECLINE) plus the falling movers. Use a clear, readable table or ranked rows — calm and professional, never alarmist, but never softened into vague positivity. OMIT only if genuinely nothing declined.\n"
    "19. Strategic Insights & Recommendations — actionable recommendations ONLY, spanning BOTH search and (when present) website analytics: SEO content/internal-linking/CTR/title improvements, target near-page-1 queries, defend declining queries/pages, plus double-down on best-converting channels and lift engagement/conversion where the analytics support it.\n"
    "20. Closing Slide — key takeaways, honest momentum summary, professional thank-you page with the reporting period in a slim footer."
)


# Structure for the GA4-only website-analytics deck (Google Analytics on-site behaviour).
# No Search Console / organic-search slides — this deck is purely GA4.
GA4_STRUCTURE = (
    "1. Cover Slide — site/property name, 'Website Analytics Report', and the REPORTING PERIOD date range as the subtitle, on a professional hero visual. NO KPI numbers/metric chips on the cover — keep it clean.\n"
    "2. Executive Summary — CHART-LED slide: one or two large Plotly charts (sessions/users trend or channel mix) as the focus, with the high-level KPIs as a slim compact strip along the top or bottom (NOT large number cards with empty space). Strongest positives, short strategic summary.\n"
    "3. Audience & Engagement — from AUDIENCE & ENGAGEMENT: sessions, total users, new users, pageviews, engagement rate, bounce rate, avg session duration and conversions with their period-over-period change; highlight key wins visually.\n"
    "4. Sessions Over Time — from SESSIONS OVER TIME: the daily series as the hero chart, an area/line of sessions & users (with conversions overlaid if present). Include a VISIBLE LEGEND naming each series and clear axis titles.\n"
    "5. Traffic by Channel (REQUIRED — always include this slide) — from TRAFFIC BY CHANNEL: how sessions/users/conversions split across channels (organic, direct, paid, referral, social, etc.). Use a horizontal bar or donut of sessions by channel paired with a small labelled table; highlight the strongest acquisition channels and frame conversions positively.\n"
    "6. Sessions by Device — from BY DEVICE: a donut or bar of sessions by desktop/mobile/tablet with each device's session share; note the period-over-period movement. OMIT this slide if BY DEVICE is (none).\n"
    "7. Geographic Distribution — from GEOGRAPHY: a Plotly CHOROPLETH world map (\"type\":\"choropleth\", \"locationmode\":\"country names\") shaded by sessions, paired with a top-countries bar. Keep the geo clean (accent colourscale, transparent bg, no coastline/frame chartjunk).\n"
    "8. Needs Attention — What's Dropping (REQUIRED — always include this slide) — every declining metric, channel, device or geography RANKED BY IMPACT, each with its real movement (previous → current, % change), the likely cause, and the specific fix. Calm and professional, never alarmist, but never softened into vague positivity. OMIT only if genuinely nothing declined.\n"
    "9. Strategic Insights & Recommendations — actionable recommendations ONLY, grounded in the analytics: double-down on best-converting channels, lift engagement/reduce bounce where the data supports it, grow the strongest audiences/geographies, and improve conversion paths.\n"
    "10. Closing Slide — key takeaways, honest momentum summary, professional thank-you page with the reporting period in a slim footer."
)


# Structure for the Bing (Microsoft) organic-search deck. Bing gives clicks/impressions +
# top queries/pages; CTR and period deltas are derived upstream. The AI Search Visibility
# slide (Copilot citations) is OMITTABLE — only present when the user uploads the AI
# Performance CSV export (Bing has no AI API yet).
BING_STRUCTURE = (
    "1. Cover Slide — site/domain, 'Bing Search Report', and the REPORTING PERIOD date range as the subtitle, on a professional hero visual. NO KPI numbers/metric chips on the cover — keep it clean.\n"
    "2. Executive Summary — CHART-LED slide: one or two large Plotly charts (clicks/impressions trend) as the focus, with the high-level KPIs as a slim compact strip along the top or bottom (clicks, impressions, CTR — and 'Total AI Citations' too WHEN the AI SEARCH VISIBILITY section is present). NOT large number cards with empty space. Strongest positives, short strategic summary.\n"
    "3. Bing Search Performance — clicks, impressions, CTR with period-over-period change; highlight key wins visually. Note this is Microsoft Bing organic search (distinct from Google).\n"
    "4. Performance Over Time — from PERFORMANCE OVER TIME: a Plotly chart of daily clicks & impressions (filled AREA or combo). Include a VISIBLE LEGEND naming each series and clear axis titles.\n"
    "5. Top Queries (REQUIRED — always include this slide) — a clean, readable TABLE of the top ~10 queries by clicks with columns Query · Clicks · Impressions · CTR · Avg position. This is a real labelled table so the boss can read exactly which keywords drive Bing traffic.\n"
    "6. Top Pages (REQUIRED — always include this slide) — a readable TABLE of the best-performing landing pages by clicks (Page · Clicks · Impressions · CTR); shorten long URLs to their path. Focus on strengths.\n"
    "7. AI Search Visibility (OMIT this ENTIRE slide if there is no AI SEARCH VISIBILITY section in the data) — the 'reporting on AI' slide: show Total AI Citations and Average Cited Pages as KPIs, make the CITATIONS OVER TIME daily series the hero chart (a Plotly filled AREA/line of citations), and call out the peak citation day. Frame it as how often this site is cited as a source in Microsoft Copilot / Bing AI-generated answers — a growing AI-search visibility signal.\n"
    "8. Needs Attention — What's Dropping (REQUIRED — always include this slide) — every declining query, landing page or metric RANKED BY IMPACT, each with its real movement (previous → current, % change), the likely cause, and the specific fix. Calm and professional, never alarmist, but never softened into vague positivity. OMIT only if genuinely nothing declined.\n"
    "9. Strategic Insights & Recommendations — actionable Bing-specific recommendations ONLY: grow Bing visibility, target near-page-1 Bing queries, improve CTR on high-impression queries, and (when AI data present) sustain/grow AI-citation momentum by strengthening the most-cited content. Do NOT reference Google.\n"
    "10. Closing Slide — key takeaways, honest momentum summary, professional thank-you page with the reporting period in a slim footer."
)


# ---------------------------------------------------------------------------
# DESIGN SYSTEM — always appended (like HTML_CONTRACT) so EVERY deck, including
# custom prompts, gets the same template-grade consistency. This is what closes the
# gap to polished tools: cohesive tokens + a fixed component/layout vocabulary so all
# slides look like one designed template instead of N improvised ones.
# ---------------------------------------------------------------------------
DESIGN_SYSTEM = """=== DESIGN SYSTEM (required — build the deck like a senior studio template) ===
Define a single design language ONCE in :root and reuse it on every slide. Do NOT invent
per-slide colours, fonts or spacing — everything references the tokens below.

1. DESIGN TOKENS — put this in the <style> and use var(...) everywhere:
   :root{
     --bg: <chosen background>;        /* page/canvas ground */
     --surface: <card surface>;        /* cards, panels */
     --ink: <primary text>;            /* headings/body — must be high-contrast on --bg */
     --muted: <secondary text>;        /* captions, axis labels */
     --accent: <ONE vivid accent>;     /* the single pop colour */
     --accent-2: <supporting tone>;    /* second series / subtle fills */
     --line: <hairline/border colour>;
     --font-display: '<display font>', serif;   /* headlines */
     --font-body: '<body font>', sans-serif;    /* everything else */
   }
   Type scale (use consistently — this is a 1920x1080 canvas, so type is LARGE):
   display 96-144px, h2 52-64px, kpi-number 64-84px, body 26-30px, caption 18-20px.
   Spacing in multiples of 12px. Border-radius consistent.

2. SHARED COMPONENTS — define and reuse these classes (consistent on every slide).

   THE SLIDE CHROME (this is what makes the deck read as ONE designed system — every content
   slide carries the header, and most carry the takeaway bar and footer):
   - .slide-header : the header block = .eyebrow + h2 title + optional .subtitle + .rule.
   - .eyebrow      : small ALL-CAPS letter-spaced kicker above the title (e.g. "BUDGET EFFICIENCY REVIEW").
   - .subtitle     : ONE quiet line under the title (--muted, or --accent for emphasis).
   - .rule         : a thin FULL-WIDTH accent hairline closing the header off from the content.
   - .takeaway     : a FULL-WIDTH SOLID DARK band pinned near the bottom, containing
                     .takeaway-label (small ALL-CAPS, in --accent) + the takeaway sentence in white.
                     This carries the slide's key takeaway and is the deck's recurring motif.
   - .footer       : a thin muted row at the very bottom — left: "<client> — <report name>",
                     right: the period/date (and "Confidential" where appropriate).

   CONTENT COMPONENTS:
   - .kpi-tile  : a SOLID TINT tile = small label (--muted) + BIG number (--font-display) + a .delta
                  chip. Variants .tile-dark (inverted, white text) and .tile-accent (solid accent/green)
                  to spotlight ONE hero metric.
   - .delta     : a small solid pill showing movement — .delta-good (green) for a win,
                  .delta-bad (red) for a decline, .delta-warn (amber) for caution. ALWAYS semantic.
   - .panel     : a SOLID tint block (--tint or --tint-2) with generous padding — the default grouping
                  device. .panel-dark for an inverted panel. NO gradients, NO single-side accent stripes.
   - .stat-big  : a very large hero number (with its label) used inside a panel — the "74.9%" moment.
   - table      : dark header row (--dark bg, white text), zebra body rows (--tint-2 on alternate rows),
                  right-aligned numerics, and CELL VALUES COLOURED SEMANTICALLY (a declining value is
                  red, a winning value green). A final total row may be inverted dark.
   - .chip      : small pill label for tags.
   - .pageno    : small corner index "03 / 09".

3. THE SLIDE SKELETON (every content slide, no exceptions):
   Make every .slide a full-height column: display:flex; flex-direction:column; height:1080px;
   padding:84px 108px.
     .slide-header (eyebrow + title + subtitle + rule)
     MAIN content region — flex:1; min-height:0 — this GROWS to fill all remaining height
     .takeaway  (the dark band, on content slides that have a key takeaway)
     .footer
   The main region's chart/panels/table MUST stretch to consume that space — size charts large
   (~640-780px tall), enlarge panels and spacing. No slide may leave more than ~10% empty vertical
   space. Cover / section / closing slides are the exception: they are posters and skip the chrome.

4. LAYOUT ARCHETYPES — EVERY <section class="slide"> must ALSO carry one archetype class and follow
   its structure. The archetype describes the MAIN CONTENT REGION (the chrome above stays constant).
   - .layout-cover      : title page — a two-column split ~60/40: the LEFT is a clean panel holding
                          the eyebrow report label, a very large bold title, a short accent rule, the
                          client name, then quiet metadata lines (period, account, "Prepared by").
                          The RIGHT is a full-height hero PHOTO (<img class="ai-img">, object-cover)
                          with a thin accent stripe at the seam. NO KPI numbers on the cover.
   - .layout-section    : a chapter divider — huge number + section title on a solid saturated/dark field.
   - .layout-kpi-tiles  : a grid of .kpi-tile (4 across, 1-2 rows), each label + big number + .delta
                          chip; make ONE tile .tile-dark or .tile-accent to spotlight the hero metric.
   - .layout-kpi-strip  : a slim row of 3-5 metrics separated by thin vertical dividers (no boxes) ABOVE
                          ONE large chart that fills the rest of the height.
   - .layout-chart-rail : the workhorse — ONE large chart taking ~70% width, plus a right rail of a
                          .panel callout (bullets or a .stat-big) and an optional quiet note panel.
   - .layout-chart-table: a chart on one side and a data TABLE on the other (dark header, zebra rows,
                          semantic cell colours) — use for before/after and ranked breakdowns.
   - .layout-stat-panel : a solid DARK or ACCENT panel on one side carrying a giant number/short claim
                          + 2-3 mini-stats, and a chart (plus a small panel) on the other.
   - .layout-two-chart  : two charts side by side, each with its own small heading, plus a row of 2-4
                          tint stat panels beneath them.
   - .layout-panels     : 2-3 equal SOLID tint panels side by side, each a heading + bullet list
                          (e.g. Strengths / Opportunities / Observations).
   - .layout-priority   : ranked full-width rows — each row = a solid priority chip on the left
                          ("Priority 1 / High Impact"), the recommendation title + one-line rationale
                          in the middle, and a bordered "Expected Impact" box on the right. Tint the
                          top-priority row with --tint.
   - .layout-phases     : a phased roadmap — a thin timeline of dots/connectors across the top, then
                          3 equal tint columns (number + phase title + bullets) each with an "Outcome"
                          box pinned at the column's bottom.
   - .layout-split      : the generic two-column fallback ~40/60, both stretched full height — text/
                          metrics on one side, a large chart, table or full-height photo on the other.
   - .layout-list       : icon bullets spread over the full height — .ai-icon + bold lead-in + one line.
   - .layout-comparison : two/three panels compared side by side, full height.
   - .layout-quote      : one large pull-quote / single headline insight, vertically centred.
   - .layout-closing    : closing poster — big type on a solid field + the period in a slim footer.
   Vary archetypes across the deck; NEVER repeat the same layout on consecutive content slides."""


# A short menu of vetted, cohesive palette + Google-font pairings. The model picks ONE
# that fits the industry and binds it to the tokens — prevents clashing colour/type.
THEME_PRESETS = """=== THEME PRESETS (pick exactly ONE that best fits THIS site's industry, then bind it to the tokens) ===
Choose the preset whose mood matches the brand; two different industries should rarely land on the same one.
A. "Editorial Cream" — bg #FAF7F0, surface #FFFFFF, ink #1A1A1A, muted #6B6B6B, accent #E4572E, accent-2 #2A4D69; display 'Fraunces', body 'Inter'. (lifestyle, retail, hospitality)
B. "Modern Mono" — bg #0E0E10 (dark) for covers / #F5F5F4 content, ink #111 / #FAFAFA, accent #FACC15, accent-2 #64748B; display 'Space Grotesk', body 'Inter'. (tech, SaaS, startups)
C. "Clean Corporate" — bg #FFFFFF, surface #F8FAFC, ink #0F172A, muted #64748B, accent #2563EB, accent-2 #0EA5E9; display 'Archivo', body 'Inter'. (B2B, finance, consulting)
D. "Warm Premium" — bg #FBF6F1, surface #FFFFFF, ink #20140E, muted #7A6A5F, accent #B45309, accent-2 #166534; display 'Playfair Display', body 'Source Sans 3'. (luxury, craft, food & beverage)
E. "Bold Vermillion" — bg #FFF8F4, surface #FFFFFF, ink #1C1917, muted #78716C, accent #DC2626, accent-2 #1E293B; display 'Bricolage Grotesque', body 'Inter'. (media, sport, bold consumer brands)
F. "Forest Botanical" — bg #F6F4EC, surface #FFFFFF, ink #14241B, muted #5E6B60, accent #2F7A4D, accent-2 #B07B2C; display 'Fraunces', body 'Inter'. (health, wellness, sustainability, agriculture)
G. "Deep Teal" — bg #F2F6F6, surface #FFFFFF, ink #07292B, muted #5A7173, accent #0E8388, accent-2 #E2A33B; display 'Syne', body 'Inter'. (clinics, dental, science, marine)
H. "Royal Plum" — bg #FAF6FB, surface #FFFFFF, ink #1E1024, muted #6E5C72, accent #7C3AED, accent-2 #DB2777; display 'Instrument Serif', body 'Inter'. (beauty, fashion, creative agencies)
I. "Ink & Gold" — bg #0F1115 (dark) for covers / #F4F1EA content, ink #0F1115 / #F7F4EC, accent #C9A227, accent-2 #8C7A4B; display 'Libre Caslon Display', body 'Source Sans 3'. (legal, jewellery, premium services)
J. "Slate Industrial" — bg #F1F2F4, surface #FFFFFF, ink #15191E, muted #5B6470, accent #EA580C, accent-2 #334155; display 'Archivo', body 'Inter'. (manufacturing, automotive, logistics, construction)
K. "Coastal Sky" — bg #F3F7FB, surface #FFFFFF, ink #0B2540, muted #5C7287, accent #0EA5E9, accent-2 #14B8A6; display 'Space Grotesk', body 'Inter'. (travel, real estate, education)
L. "Berry Cream" — bg #FBF5F4, surface #FFFFFF, ink #2A1116, muted #7A5E60, accent #BE123C, accent-2 #2563EB; display 'Playfair Display', body 'Inter'. (restaurants, events, florists, boutique retail)
Load the chosen fonts from fonts.googleapis.com. You MAY tune the accent toward the brand's industry, but keep the palette cohesive and high-contrast."""


# One complete worked example anchors the model far more than instructions alone. Generic
# placeholder content only (no numbers that could be mistaken for real data).
DESIGN_EXEMPLARS = """=== EXEMPLAR (the SKELETON + chrome to match — do NOT copy its words, theme or composition) ===
This shows the agency-grade slide SYSTEM: every content slide = header (eyebrow + CLAIM title +
subtitle + rule) -> a main region that fills all remaining height -> the dark EXECUTIVE TAKEAWAY band
-> the muted footer. Follow the ART DIRECTION + the planned LAYOUT for this deck's actual composition:
<section class="slide layout-chart-rail" style="display:flex;flex-direction:column;height:1080px;padding:84px 108px">
  <div class="slide-header">
    <span class="eyebrow">Budget Efficiency Review</span>
    <h2>Headroom Beyond the Budget Cap</h2>
    <p class="subtitle">Utilization and incremental conversion potential</p>
    <div class="rule"></div>
  </div>
  <div style="flex:1;min-height:0;display:grid;grid-template-columns:7fr 3fr;gap:36px;margin:32px 0">
    <div id="chart3" style="width:100%;height:100%"></div>
    <div style="display:flex;flex-direction:column;gap:24px">
      <div class="panel"><h3>Reading the Trend</h3><ul><li>—</li><li>—</li></ul></div>
      <div class="panel panel-quiet"><span class="stat-big">—</span><span>label</span></div>
    </div>
  </div>
  <div class="takeaway"><span class="takeaway-label">Executive Takeaway</span><p>The one sentence the client should remember.</p></div>
  <div class="footer"><span>Client — Report name</span><span>Period</span></div>
  <script type="application/json" class="plotly-spec" data-target="chart3">{"data":[{"type":"bar","x":["Mar","Apr"],"y":[3,5]}],"layout":{}}</script>
</section>
And the cover — a two-column split (NOT a full-bleed photo): a clean panel of type on one side, a
full-height photo on the other, with a thin accent stripe at the seam. NO metric chips on the cover:
<section class="slide layout-cover" style="display:grid;grid-template-columns:6fr 4fr;height:1080px">
  <div style="display:flex;flex-direction:column;justify-content:center;padding:96px 108px">
    <span class="eyebrow">Google Ads Performance Review</span>
    <h1>Performance Audit Report</h1>
    <div class="rule" style="width:120px"></div>
    <p class="client">Client Name</p>
    <p class="meta">15 May – 15 June 2026 · vs previous period</p>
    <p class="meta">Prepared by TBS Marketing</p>
  </div>
  <img class="ai-img" data-prompt="premium product context photo, editorial lighting, no text" style="width:100%;height:100%;object-fit:cover">
</section>
Every real slide must fill all 1920x1080 (no bottom band), use the tokens/components above, put a
relevant icon on every KPI and list bullet, and use ai-img photos liberally — ALWAYS a full-bleed
hero photo on the cover, and photo backgrounds/side-panels on most content slides."""


_PRESET_LETTERS = "ABCDEFGHIJKL"  # the 12 THEME_PRESETS entries

# TBS Marketing brand palette (from the logo — blue / green / amber). Used as the DEFAULT deck
# colour/look so decks carry the agency's house identity. _apply_theme forces the accents onto the
# rendered HTML regardless of the model; _TBS_STYLE_DIRECTIVE fixes the fonts + ground.
TBS_PALETTE = {
    "accent": "#3C8DD9",    # TBS blue (primary pop / data emphasis)
    "accent2": "#79B84B",   # TBS green (wins / secondary series)
    "accent3": "#F4B740",   # TBS amber (caution)
    "bad": "#C4553B",       # terracotta red — declines / losing ground
    "ink": "#1F2937",       # near-black navy ink
    "dark": "#23262B",      # solid dark panel (inverted tiles, table headers, takeaway bar)
    "muted": "#6E7075",
    "tint": "#E8F4FB",      # solid light-blue tint panel
    "tint2": "#F2F3EC",     # solid light-neutral tint panel
    "bg": "#FFFFFF",        # clean white ground
    "surface": "#FFFFFF",
}
# A BOLD GEOMETRIC SANS is correct here — the agency reference decks are exactly this on white and
# read as premium. (An earlier switch to a display serif on cream was a misdiagnosis: the flatness
# came from missing slide chrome + flat compositions, not from the typeface.)
TBS_FONTS = {"display": "Poppins", "body": "Inter"}

_TBS_STYLE_DIRECTIVE = (
    "=== ASSIGNED HOUSE STYLE — TBS MARKETING (use THIS exactly) ===\n"
    "A modern, confident CONSULTING house style: clean white ground, bold geometric type, solid tint "
    "panels, and colour used to carry MEANING. Never a generic corporate slide template.\n"
    f"- GROUND: --bg {TBS_PALETTE['bg']} (white), --surface {TBS_PALETTE['surface']}, "
    f"--ink {TBS_PALETTE['ink']}, --muted {TBS_PALETTE['muted']}.\n"
    f"- PANELS are SOLID tints — never gradients: --tint {TBS_PALETTE['tint']} (light blue), "
    f"--tint-2 {TBS_PALETTE['tint2']} (light neutral), and --dark {TBS_PALETTE['dark']} for inverted "
    "panels/table headers/the takeaway bar. Cover, section dividers and closing may use ONE full-bleed "
    "saturated field in the accent, the green or the dark.\n"
    f"- TYPE: load '{TBS_FONTS['display']}' as --font-display and '{TBS_FONTS['body']}' as --font-body "
    "from fonts.googleapis.com. Titles are LARGE and BOLD (clamp ~64-104px), tight line-height, sentence "
    "case; a small ALL-CAPS letter-spaced kicker sits above the title; the subtitle is one quiet line.\n"
    f"- SEMANTIC COLOUR (non-negotiable): wins/positive -> {TBS_PALETTE['accent2']} (green); "
    f"declines/at-risk -> {TBS_PALETTE['bad']} (terracotta red); caution -> {TBS_PALETTE['accent3']} "
    f"(amber); neutral emphasis/brand -> {TBS_PALETTE['accent']} (blue). Deltas, table cells and chart "
    "series MUST take these colours. NEVER colour a decline green.\n"
    "(The --accent / --accent-2 are set automatically to the TBS brand colours — just use this ground, "
    "these fonts and this system.) Do not substitute a different preset or background."
)


def _seeded_preset(seed: str) -> str:
    """Deterministically map a site (domain) to ONE theme preset letter, so the same
    client always gets the same typographic identity and different clients spread across
    the menu — instead of the model defaulting most decks to the first preset. Uses a
    stable hash (not Python's salted hash) so it's consistent across processes/restarts."""
    import hashlib
    h = int(hashlib.sha1(seed.strip().lower().encode()).hexdigest(), 16)
    return _PRESET_LETTERS[h % len(_PRESET_LETTERS)]


def _seed_directive(seed: str) -> str:
    """A directive that pins the deck to the seed's preset for background + fonts + mood.
    The accent is overridden downstream (_apply_theme) to the real brand colour, so we only
    fix the typographic/background identity here — that's what makes clients look distinct."""
    letter = _seeded_preset(seed)
    return (
        "=== ASSIGNED THEME PRESET (use THIS one) ===\n"
        f"For THIS site, use THEME PRESET {letter} — adopt its BACKGROUND, SURFACE, font "
        f"pairing (display + body) and overall mood exactly. (The --accent / --accent-2 will be "
        f"set to the brand's own colour automatically, so don't worry about matching the preset's "
        f"accent — just take its grounds, fonts and feel.) Do not substitute a different preset."
    )


def _preset_directive(letter: str) -> str:
    """Force a specific THEME_PRESET letter (user picked a named style)."""
    letter = (letter or "A").upper()
    if letter not in _PRESET_LETTERS:
        letter = "A"
    return (
        "=== ASSIGNED THEME PRESET (use THIS one) ===\n"
        f"Use THEME PRESET {letter} — adopt its BACKGROUND, SURFACE, font pairing (display + body) and "
        "overall mood exactly. (The --accent / --accent-2 are set to the chosen brand colour "
        "automatically, so just take its grounds, fonts and feel.) Do not substitute a different preset."
    )


def _style_directive(style: Optional[str], seed: Optional[str]) -> str:
    """Pick the fonts/background directive: 'tbs' → the fixed TBS house style; 'auto' (or empty) →
    the per-domain preset; a single letter A–L → that specific preset."""
    s = (style or "tbs").strip()
    if s == "tbs":
        return _TBS_STYLE_DIRECTIVE
    if s == "auto":
        return _seed_directive(seed or "")
    if len(s) == 1 and s.upper() in _PRESET_LETTERS:
        return _preset_directive(s)
    return _TBS_STYLE_DIRECTIVE


# Art-direction axes rotated per-seed so different clients get visibly different decks — the single
# worked EXEMPLAR otherwise anchors every deck to the same cover + kpi-strip rhythm.
_COVER_STYLES = [
    "a FULL-BLEED photographic cover (ai-img) with a dark gradient and an oversized title",
    "a BOLD single saturated colour-field poster cover — huge type, NO photo, generous negative space",
    "a TYPOGRAPHIC cover: a giant display headline on the ground with a thin ruled/dot grid and a small photo inset",
    "a SPLIT cover — a full-height photo on one side, a colour field carrying the title on the other",
]
_LEAD_LAYOUTS = [
    "a kpi-strip (a slim KPI row above ONE hero chart)",
    "a split 40/60 (text + KPIs on one side, a large chart on the other)",
    "a comparison of 2-3 cards",
    "a two-chart data spread side by side",
]
_MOTIFS = [
    "a thin dot-grid",
    "a slightly rotated outlined frame",
    "pill / tape-style labels",
    "a hairline rule paired with index numbers",
    "a single small filled star or circle accent",
]
_HEADLINE_STYLES = [
    "oversized high-contrast SERIF headlines with ONE italic accent word",
    "a bold grotesque with tiny ALL-CAPS letter-spaced kickers over large sans headlines",
    "editorial serif headlines with wide letter-spaced eyebrows and lots of air",
]


def _variant_directive(seed: str) -> str:
    """Deterministically pick an art-direction 'recipe' from the seed so two different sites produce
    visibly different decks (cover treatment, opening layout rhythm, motif, headline style) instead of
    every deck copying the single EXEMPLAR's cover+kpi-strip look. Per-seed ⇒ a given client stays
    consistent while clients spread across the axes."""
    import hashlib
    h = int(hashlib.sha1(("v:" + (seed or "").strip().lower()).encode()).hexdigest(), 16)
    cover = _COVER_STYLES[h % len(_COVER_STYLES)]
    lead = _LEAD_LAYOUTS[(h // 7) % len(_LEAD_LAYOUTS)]
    motif = _MOTIFS[(h // 13) % len(_MOTIFS)]
    head = _HEADLINE_STYLES[(h // 17) % len(_HEADLINE_STYLES)]
    return (
        "=== ART DIRECTION FOR THIS DECK (make it visibly DIFFERENT from other decks) ===\n"
        f"- COVER: use {cover}.\n"
        f"- OPENING RHYTHM: lead your first data slide with {lead} (do NOT default to the exemplar's "
        "cover→kpi-strip order).\n"
        f"- SIGNATURE MOTIF: carry {motif} across slides (never over charts/tables).\n"
        f"- HEADLINES: {head}.\n"
        "Treat these as the deck's identity and vary layouts slide to slide — the EXEMPLAR below is a "
        "quality/height reference ONLY, not a layout to copy."
    )


def _structure_directive(creativity: str, structure: str) -> str:
    """How strictly the model must follow the structure. The `structure` string is the same
    set of topics/data at every level — only the freedom to reshape it into slides changes.
    All levels still keep the HTML contract, chart mechanism and validation guardrails."""
    themes = structure or DEFAULT_STRUCTURE
    if creativity == "structured":
        return (
            "SLIDE SEQUENCE — produce ONE slide for EACH numbered item below, IN ORDER. This is the "
            "authoritative outline; follow it faithfully:\n\n" + themes + "\n\n"
            "SLIDE-COUNT DISCIPLINE (critical): a SEPARATE slide for every applicable item — do NOT "
            "merge items onto one slide, do NOT shorten to a generic 8-10 slide summary, and NEVER "
            "skip an item marked REQUIRED (especially the keyword-opportunity bubble and the "
            "Needs Attention / what's-dropping slide). The ONLY items you may omit are those explicitly "
            "marked \"OMIT … if …\" whose data is absent. Expect a rich deck of ~12-20 slides depending "
            "on which optional data is present.")
    if creativity == "creative":
        return (
            "DESIGN THIS DECK YOUR OWN WAY. The list below is a CHECKLIST OF WHAT TO COVER, not a slide "
            "list. YOU decide the number of slides and their order, and invent a DISTINCT layout and "
            "visual treatment for each slide — be bold and editorial so no two slides look alike. You "
            "may split, merge, reorder, or add connective/insight slides for the strongest narrative. "
            "Cover every REQUIRED item and use ONLY the real data. Prioritise sharp takeaways and "
            "concrete recommendations over restating numbers.\n\nCOVER:\n" + themes)
    # balanced (default)
    return (
        "The list below is the set of THEMES to cover. Cover every REQUIRED item and all key data, but "
        "CHOOSE YOUR OWN slide count and order and give each slide a distinct layout — you may split, "
        "merge, or resequence for the best narrative flow (aim ~8-14 slides). Favour clear insights and "
        "recommendations, not just restating numbers.\n\nTHEMES TO COVER:\n" + themes)


def build_prompt(data_brief: str, *, prompt: Optional[str] = None,
                 brand: Optional[str] = None, structure: Optional[str] = None,
                 seed: Optional[str] = None, creativity: str = "balanced",
                 variant_seed: Optional[str] = None, style: Optional[str] = "tbs") -> str:
    """Fill the (possibly user-customised) prompt template with brand/structure/data,
    then append the rendering contract + design system + theme presets + exemplar so
    every deck — default or custom — gets the same template-grade quality bar.

    `creativity` (structured|balanced|creative) sets how freely the model may reshape the
    structure into slides. If `seed` (e.g. the site domain) is given, a deterministic preset
    is pinned so the same client always gets the same typographic identity."""
    template = prompt or DEFAULT_DECK_PROMPT
    directive = _structure_directive(creativity, structure or DEFAULT_STRUCTURE)
    # Use replace (not str.format): prompts contain literal CSS braces.
    filled = (
        template
        .replace("{brand}", brand or DEFAULT_BRAND)
        .replace("{structure_directive}", directive)
        .replace("{structure}", structure or DEFAULT_STRUCTURE)  # back-compat for custom prompts
        .replace("{data}", data_brief)
    )
    # If a custom prompt forgot the placeholders, append them anyway so the model still gets the
    # structure directive and the data (otherwise they're silently dropped).
    if "{structure_directive}" not in template and "{structure}" not in template:
        filled = filled + "\n\n" + directive
    if "{data}" not in template:
        filled = filled + "\n\nDATA (use ONLY this):\n" + data_brief
    parts = [filled, HTML_CONTRACT, DESIGN_SYSTEM, THEME_PRESETS]
    parts.append(_style_directive(style, seed))      # fonts/ground: TBS house / per-site / a preset
    # Art direction varies per GENERATION (variant_seed), so two models / re-runs of the SAME site
    # get different compositions — pinning it to the domain made every deck for a site look identical.
    parts.append(_variant_directive(variant_seed or seed or ""))
    parts.append(DESIGN_EXEMPLARS)
    return "\n\n".join(parts)


def _clean_html(text: str) -> str:
    """Strip any stray markdown fences the model may wrap around the HTML."""
    t = text.strip()
    if t.startswith("```"):
        t = t.split("\n", 1)[1] if "\n" in t else t
        if t.endswith("```"):
            t = t[: t.rfind("```")]
    i = t.lower().find("<!doctype")
    if i == -1:
        i = t.lower().find("<html")
    return t[i:].strip() if i != -1 else t.strip()


_SCRIPT_STYLE_RE = re.compile(r'<(script|style)\b[^>]*>.*?</\1>', re.IGNORECASE | re.DOTALL)


def _match_brace(s: str, start: int) -> int:
    """Index of the '}' that closes the '{' at `start`, respecting JSON strings. -1 if none."""
    depth, in_str, esc = 0, False, False
    for k in range(start, len(s)):
        c = s[k]
        if in_str:
            if esc:
                esc = False
            elif c == "\\":
                esc = True
            elif c == '"':
                in_str = False
        elif c == '"':
            in_str = True
        elif c == "{":
            depth += 1
        elif c == "}":
            depth -= 1
            if depth == 0:
                return k
    return -1


def _strip_blobs(s: str, token: str) -> str:
    """Delete every balanced {...} object that begins at `token` and carries BOTH a "data"
    and a "layout" key (a leaked Plotly chart spec). Any other inline JSON is preserved."""
    out, i, n = [], 0, len(s)
    while i < n:
        j = s.find(token, i)
        if j == -1:
            out.append(s[i:])
            break
        end = _match_brace(s, j)
        if end == -1:
            out.append(s[i:])
            break
        blob = s[j:end + 1]
        if '"data"' in blob and '"layout"' in blob:   # a leaked chart spec → drop it
            out.append(s[i:j])
        else:                                          # other inline JSON → keep it
            out.append(s[i:end + 1])
        i = end + 1
    return "".join(out)


def _strip_leaked_specs(html: str) -> str:
    """Remove any chart-spec JSON the model printed as VISIBLE text (a bare object carrying
    both "data" and "layout", in any key order). Legit specs live inside <script class=
    "plotly-spec"> and must be preserved, so script/style bodies are masked out before
    scanning. Cheap first net; the render-time DOM cleanup is the authoritative guarantee
    (it also catches HTML-escaped/oddly-wrapped leaks this regex pass can't see)."""
    if '"data"' not in html or '"layout"' not in html:
        return html
    blocks: list = []

    def _hide(m):
        blocks.append(m.group(0))
        return f"\x00B{len(blocks) - 1}\x00"

    masked = _SCRIPT_STYLE_RE.sub(_hide, html)
    # Catch both data-first and layout-first objects (with or without a space after '{').
    for token in ('{"data"', '{ "data"', '{"layout"', '{ "layout"'):
        masked = _strip_blobs(masked, token)
    return re.sub(r"\x00B(\d+)\x00", lambda m: blocks[int(m.group(1))], masked)


_DECK_SYSTEM_PROMPT = "You are an award-winning presentation designer who outputs only clean, self-contained HTML."


def _build_repair_prompt(html: str, instructions: str) -> str:
    """A focused follow-up asking the model to fix ONLY the listed problems and
    re-emit the whole document. Keeps the deck self-contained per HTML_CONTRACT."""
    return (
        "The HTML presentation below has specific problems that must be fixed. Fix ONLY "
        "the problems listed; preserve all other content, wording, numbers, layout and "
        "design exactly as they are.\n\n"
        "PROBLEMS TO FIX:\n" + instructions + "\n\n"
        "Return the COMPLETE corrected HTML document and NOTHING ELSE — no markdown, no "
        "commentary. It must still satisfy this contract:\n\n" + HTML_CONTRACT + "\n\n"
        "=== CURRENT HTML ===\n" + html
    )


def _make_image_prewarmer(image_cache: Dict[str, "asyncio.Task"]):
    """Build an on_delta callback that, as the deck streams in, starts generating an
    image the moment each <img class="ai-img" data-prompt="..."> placeholder is fully
    written — so photos render concurrently with the rest of the slide-writing instead
    of serially afterward. Tasks are keyed by prompt in image_cache for resolve_ai_images
    to await. Bounded by MAX_DECK_IMAGES and _IMG_CONCURRENCY (rate-limit safe)."""
    from services.image_service import generate_image
    sem = asyncio.Semaphore(_IMG_CONCURRENCY)

    async def _gen(prompt: str):
        async with sem:
            return await generate_image(prompt)

    async def on_delta(buf: str):
        if len(image_cache) >= MAX_DECK_IMAGES or "ai-img" not in buf:
            return
        for m in _AI_IMG_RE.finditer(buf):
            if len(image_cache) >= MAX_DECK_IMAGES:
                break
            p = _img_prompt(m.group(0))
            if p and p not in image_cache:
                image_cache[p] = asyncio.create_task(_gen(p))

    return on_delta


# ---------------------------------------------------------------------------
# 3-LAYER PIPELINE (opt-in): plan → per-slide insights → HTML. Each layer can use
# a different model, so the same deck can be A/B'd across models per stage. The final
# (HTML) layer still emits the same ai-img/plotly-spec contract, so image prewarming,
# keyword-bubble enforcement and leaked-spec stripping are unchanged.
# ---------------------------------------------------------------------------
_GUIDELINES_SYSTEM = ("You are a design director. You output a concise, concrete design brief in markdown "
                      "— no preamble, no code fences.")

_GUIDELINES_PROMPT = """Write the COMMON DESIGN GUIDELINES for ONE client presentation, derived from the
DATA below. Every slide of this deck is designed by a SEPARATE agent from this brief alone, so it must be
concrete, prescriptive and unambiguous. Output markdown only.

Cover exactly these sections:

## Client & Context
The client/site, the audience (a busy business owner / executive — NON-technical: translate platform
metrics into business outcomes), the tone (strategic, consultative, authoritative and HONEST — diagnose
problems but always pair them with the fix and its expected impact), and the canvas: {canvas}px, no
scrolling.

## Colour Palette & Semantics
Take the ground / surface / ink / fonts from the assigned HOUSE STYLE below. Use these EXACT brand accents:
--accent {accent} (brand + neutral data emphasis) and --accent-2 {accent2} (secondary series).
Then FIX the semantic colours (state the exact hex for each) and use them consistently on EVERY slide:
- GOOD / win / positive movement -> a green
- BAD / critical / losing ground -> a red
- CAUTION / needs attention -> an amber
- NEUTRAL highlight / brand emphasis -> the accent above

## Typography
The display + body faces from the HOUSE STYLE, the type scale for this canvas, and where small uppercase
letter-spaced kickers are allowed.

## Charts
Code-rendered (Plotly) — NEVER images. State: transparent paper/plot background, the gridline colour, the
axis-label colour, ALWAYS include units in axis labels (%, x, currency), colour each series by SEMANTIC (a
declining series uses the red, a winning series the green), and annotate the peak and the low point of any
trend.

## Export-safety (STRICT — this deck is exported to PowerPoint)
- NO single-side borders (no left/top accent stripes). Use a full 4-side border OR a solid tint panel.
- NO gradient backgrounds or gradient panels (a scrim over a cover photo is the only exception).
- NO translucent/faded text and no opacity on text — use a lighter palette colour instead.
- NO nested cards (a card inside a card). Group with whitespace and headings.

## Layout Principles
Top-anchor sparse content with consistent gaps; centre only on cover/hero slides. VARY the layout every
slide — no two consecutive slides share a structure. Every non-cover slide carries a visual (chart, table,
diagram or image). A thin muted footer is allowed (left = report name, centre = client, right = slide no.).

## Recurring Motif
Pick ONE motif / emotional anchor from the REAL data — the single headline number that defines this
period's story (and its hopeful counterweight if that number is a decline). Name the number explicitly.
It should recur across the deck.

DATA (derive the client, the story and the motif from this — use ONLY real numbers):
{data}"""


_PLAN_SYSTEM = "You are a presentation strategist. You output ONLY strict JSON — no markdown, no prose."

_PLAN_PROMPT = """You are planning an executive data-presentation deck. From the coverage list and DATA
below, produce a deck PLAN as STRICT JSON ONLY (no markdown fences, no commentary).

{structure_directive}

Output JSON of EXACTLY this shape:
{"slides":[{
 "n":1,
 "title":"…",
 "objective":"what this slide must achieve for the client, in one line",
 "archetype":"layout-cover",
 "key_takeaway":"the ONE sentence the client should remember from this slide",
 "layout":"AN EXPLICIT COMPOSITION — the grid (e.g. 'two columns, 5fr 4fr'), what sits in each region, alignment, padding, and any full-bleed/scrim. Precise enough that a designer builds it without guessing.",
 "blocks":[{"type":"text|chart|table|image|diagram","placement":"which region of the layout","content":"what goes in it — the REAL figures/queries from the data","creative_brief":"how it should look and feel"}],
 "chart":{"type":"bar|line|combo|donut|scatter|choropleth|table|none","series":"exactly which data to plot","axes":"what each axis means, with units"},
 "image":{"needed":false,"prompt":"","aspect":"16:9","sizing":"object-cover"},
 "data_refs":["which figures/sections of the data this slide uses"]
}]}

Rules:
- Cover every item marked REQUIRED; use ONLY the provided data; never invent numbers.
- "archetype" MUST be one of: layout-cover, layout-section, layout-kpi-tiles, layout-kpi-strip,
  layout-chart-rail, layout-chart-table, layout-stat-panel, layout-two-chart, layout-panels,
  layout-priority, layout-phases, layout-split, layout-list, layout-comparison, layout-quote,
  layout-closing.
  Pick the one whose SHAPE matches the content (a ranked recommendation set -> layout-priority; a
  30/60/90 plan -> layout-phases; a before/after -> layout-chart-table; a single dominant number ->
  layout-stat-panel). Never use the same archetype on consecutive slides.
- Choose a sensible slide count per the directive above; order the slides for the strongest narrative.
- First slide is the cover; last slide is the closing.
- TITLES ARE CLAIMS, NOT LABELS. The title states what the data SAYS — "Flat Investment, Volatile
  Returns", "Visibility Constrained by Rank", "Where Demand Converts", "A Single-Engine Account" —
  never "Performance Overview" or "Top Queries". Put the section name in the eyebrow instead, and use
  the subtitle for the quiet descriptive line.
- LAYOUT IS YOUR JOB. Design each slide's composition explicitly in "layout" — the downstream designer
  EXECUTES your composition rather than improvising it. Vary it every slide; no two consecutive slides
  may share a structure. Obey the DESIGN GUIDELINES' layout principles and export-safety rules.
- IMAGE BUDGET: this is a DATA-LED deck. At most {max_images} images in the WHOLE deck. The cover always
  has one. Otherwise set "image":{"needed":false,...} — only use a photo where a real product/context
  shot genuinely adds meaning. NEVER an image of a chart, table, graph or diagram: those are
  code-rendered. Charts/tables are blocks of type "chart"/"table", never "image".

SEMANTIC ANALYSIS (YOU decide which lens earns a slide — add the one(s) the data genuinely supports,
and skip any that would be thin):
- TOPIC CLUSTERS: group the queries into meaningful themes (e.g. product lines, services, locations),
  and consider a slide showing each theme's clicks/impressions/avg position and which themes are
  growing vs slipping.
- SEARCH INTENT MIX: classify queries as informational / commercial / navigational / branded and
  consider a slide on how performance splits across intent.
- CONTENT-GAP THEMES: themes with high impressions but weak rank or CTR — where content is missing or
  underperforming.
Put the actual grouping you intend (the theme names and which queries belong to each) into that
slide's "data_refs" so the copy stage can build it from the real rows.

DECLINES: plan an honest "Needs Attention / what's dropping" slide whenever the data shows declines —
never bury them. Use the ANALYST FLAGS in the DATA (AT RISK / DEFEND / PAGE DECLINE) as its source.

DESIGN GUIDELINES (the brief every slide is built from — respect its palette semantics, layout
principles and export-safety rules when planning each composition):
{guidelines}

DATA (use ONLY this):
{data}"""


_NOTES_SYSTEM = ("You are a senior agency consultant preparing to present to a paying client. You output "
                 "only the speaker notes, one block per slide.")

_NOTES_PROMPT = """Write SPEAKER NOTES for a consultant delivering this deck to the client — what to SAY
while each slide is on screen.

For EACH slide, output exactly:
## Slide {n}
<2-5 sentences of spoken delivery: how to open the slide, how to read the key number aloud, what it means
for their business, and the transition to the next point. Quote the real figures. Where a slide reports a
decline, coach the presenter to state it plainly and calmly and pivot straight to the fix.>

Rules:
- Use ONLY the real numbers from the slide copy — never invent one.
- Spoken register (contractions, direct address), not written prose. No bullet lists, no markdown beyond
  the '## Slide N' headings.
- Land each slide's key takeaway explicitly.

DECK PLAN:
{plan}

SLIDE COPY (in order):
{copy}"""

_INSIGHTS_SYSTEM = "You are a senior data analyst who writes concise, grounded slide copy."

# ── Per-slide markdown (one artifact per page, like the reference tool) ──────────────────────
_SLIDE_MD_PROMPT = """Write the on-slide COPY for ONE slide of an executive deck, grounded ONLY in the DATA.

THIS SLIDE (from the deck plan):
{slide_json}

Output plain markdown ONLY (no HTML, no chart JSON, no commentary), in EXACTLY this shape:
KICKER: <tiny all-caps eyebrow, 2-4 words>
HEADLINE: <one punchy line — the INSIGHT, not the metric name>
BULLETS:
- <short insight citing the REAL number(s) from the DATA>
- <2-4 bullets total>
TAKEAWAY: <the slide's key_takeaway from the plan, sharpened — naming a real query / page / theme>
CHART: <none | one line: chart type + EXACTLY which series/rows to plot + what each axis means, with units>
PHOTO: <the plan's image prompt if image.needed is true | none>

Rules:
- Use ONLY real numbers from the DATA — never invent one.
- The plan already fixed this slide's OBJECTIVE and KEY TAKEAWAY — deliver them, don't re-invent them.
- Write to fill the plan's blocks: every block in the plan needs its real content here.
- Lead with the takeaway; the number is evidence, not the point.
- Report declines HONESTLY: state the real movement (previous -> current, % or position change), the
  likely cause and the fix. Never spin a drop into vague positivity, never omit it.
- PRIORITISE the ANALYST FLAGS in the DATA (pre-computed, grounded insights) — build this slide's
  narrative around any that apply to it.
- Cover ONLY this slide's purpose; the other planned slides cover the rest — do not duplicate them.

DECK PLAN (context only):
{plan}

DATA (use ONLY this):
{data}"""

# ── Shared stylesheet: the consistency backbone for independently-generated slides ───────────
_STYLESHEET_SYSTEM = "You are a senior design engineer. You output ONLY raw CSS — no HTML, no commentary."

_STYLESHEET_PROMPT = """Produce the SHARED STYLESHEET for ONE presentation deck. Output ONLY raw CSS —
no <style> tag, no markdown fences, no commentary.

Every slide of this deck is written INDEPENDENTLY using ONLY these classes and tokens, so the stylesheet
must be complete and self-sufficient. Define, exactly once:
- An @import url(...) from fonts.googleapis.com for the display + body fonts of the assigned HOUSE STYLE.
- Resets + page geometry:
    @page { size: 1920px 1080px; margin: 0; }
    * { box-sizing: border-box; }
    html,body { margin:0; padding:0; }
    .slide { width:1920px; height:1080px; position:relative; overflow:hidden; page-break-after:always; }
- :root tokens: --bg --surface --ink --muted --accent --accent-2 --line --font-display --font-body
  (take the ground/surface/ink/fonts from the assigned HOUSE STYLE below), PLUS --tint --tint-2 --dark
  (the solid panel grounds) and the semantic tokens --good --bad --caution using the exact hexes fixed
  in the DESIGN GUIDELINES' colour semantics.
- Helper classes for the semantics: .is-good .is-bad .is-caution (text colour) and .tint-good .tint-bad
  .tint-caution (solid tint panel backgrounds — NOT gradients, NOT single-side borders).
- The type scale for this 1920x1080 canvas (display 96-144px, h2 52-64px, .kpi-num 64-84px,
  body 26-30px, caption 18-20px) and a 12px spacing rhythm, applied via the component classes.
- The SLIDE CHROME classes: .slide-header .eyebrow .subtitle .rule (a thin full-width accent hairline),
  .takeaway (a FULL-WIDTH SOLID DARK band) + .takeaway-label (small ALL-CAPS in the accent), and
  .footer (a thin muted row, space-between).
- The content components: .kpi-tile (+ .tile-dark, .tile-accent), .delta (+ .delta-good, .delta-bad,
  .delta-warn — solid semantic pills), .panel (+ .panel-dark, and a --tint-2 variant), .stat-big,
  .chip, .pageno.
- TABLE styling: dark header row (white text), zebra body rows, right-aligned numerics, and helper
  classes so individual cells can take the semantic colours.
- EVERY layout archetype as a class that makes its slide a FULL-HEIGHT 1080px composition (no empty
  bands): .layout-cover .layout-section .layout-kpi-strip .layout-split .layout-list .layout-comparison
  .layout-roadmap .layout-quote .layout-closing — each following its structure in the DESIGN SYSTEM.
- .ai-img sizing helpers and a .photo-overlay scrim utility for text over photos (the ONLY place a
  gradient is permitted).
Make it genuinely designed — editorial, confident, generous whitespace — not a generic template.
Obey the DESIGN GUIDELINES' export-safety rules: no single-side borders, no gradient backgrounds/panels,
no translucent text (use a lighter palette colour), no nested cards."""

_SLIDE_HTML_SYSTEM = ("You are an award-winning presentation designer. You output ONE slide of clean, "
                      "self-contained HTML and nothing else.")

_SLIDE_ICON_NAMES = ("trending-up, trending-down, target, search, eye, mouse-pointer-click, users, globe, "
                     "bar-chart, line-chart, pie-chart, activity, arrow-up-right, arrow-down-right, "
                     "check-circle, alert-triangle, lightbulb, rocket, star, award, flag, calendar, clock, "
                     "map-pin, link, file-text, layers, filter, zap, dollar-sign, percent, shopping-cart, "
                     "smartphone, monitor, thumbs-up, refresh-cw, compass, megaphone, sparkles, gauge")

_SLIDE_HTML_PROMPT = """BUILD ONE slide of a presentation as HTML. Output EXACTLY ONE
<section class="slide {archetype}"> ... </section> and NOTHING ELSE — no markdown fences, no <html>,
no <head>, no <style> tag, no commentary.

Your job is to EXECUTE the composition the designer already specified — not to invent your own.

=== THE COMPOSITION TO BUILD (authoritative — follow it exactly) ===
LAYOUT: {layout}

BLOCKS (each one must appear, in its stated placement):
{blocks}

=== SLIDE COPY (the real content; refine wording lightly, NEVER invent or change a number) ===
{md}

This is slide {n} of {total}; archetype: {archetype}. Put "{index_label}" in the .pageno element.

RULES:
- Build the LAYOUT above precisely: its grid/columns and ratios, what sits in each region, the alignment,
  padding and any full-bleed/scrim. Do not substitute a different composition.
- SLIDE CHROME (mandatory on every CONTENT slide — cover/section/closing are posters and skip it):
  open with .slide-header (.eyebrow kicker + the claim title + optional .subtitle + .rule), then the
  main content region, then the .takeaway dark band carrying this slide's key takeaway
  (.takeaway-label "EXECUTIVE TAKEAWAY" + the sentence), then the .footer. This chrome is what makes
  the deck read as one designed system — do not omit or restyle it.
- Deltas/movements use a .delta chip with the SEMANTIC class (.delta-good / .delta-bad / .delta-warn).
- Use ONLY the shared stylesheet's tokens/classes (below). Do NOT emit a <style> tag and do NOT invent or
  redefine classes. Inline style="..." is allowed ONLY for per-slide geometry (grid/flex sizing, chart box
  dimensions, photo positioning).
- FILL the entire 1920x1080 canvas — a balanced, full-height composition with no empty band.
- COLOUR SEMANTICS (from the guidelines): a declining/at-risk number uses the BAD colour, a win uses the
  GOOD colour, a caution uses the AMBER, neutral emphasis uses the accent. Never colour a decline green.
- EXPORT-SAFETY (this deck is exported to PowerPoint — STRICT): no single-side borders (use a full 4-side
  border or a solid tint panel), no gradient backgrounds/panels (a cover photo scrim is the only
  exception), no translucent/faded text or opacity on text (use a lighter palette colour), no nested cards.
- CHART: if the composition calls for one, output a sized container
  <div id="chart{n}" style="width:...;height:..."></div> immediately followed by
  <script type="application/json" class="plotly-spec" data-target="chart{n}">{"data":[...],"layout":{...}}</script>
  with STRICTLY valid JSON. Transparent paper/plot bg, few gridlines, units in axis titles, series coloured
  by semantic. NEVER print chart JSON as visible text, and NEVER render a chart/table as an image.
- ICONS: <i class="ai-icon" data-icon="NAME"></i> — allowed NAMEs ONLY: {icons}
- {photo_rule}

=== DESIGN GUIDELINES (the deck-wide brief) ===
{guidelines}

{art_direction}

SHARED STYLESHEET (already in the document — these classes/tokens exist; use them):
{shared_css}"""


def _extract_json(text: str) -> Optional[dict]:
    """Pull the first balanced JSON object out of an LLM reply (tolerates code fences and
    surrounding prose). Returns the parsed dict, or None if nothing valid is found."""
    import json
    t = text.strip()
    if t.startswith("```"):
        t = t.split("\n", 1)[1] if "\n" in t else t
        if t.endswith("```"):
            t = t[: t.rfind("```")]
    start = t.find("{")
    if start == -1:
        return None
    end = _match_brace(t, start)
    if end == -1:
        return None
    try:
        obj = json.loads(t[start:end + 1])
        return obj if isinstance(obj, dict) else None
    except Exception:
        return None


async def _call_llm(full_prompt: str, *, system_prompt: str, provider: str,
                    on_progress: ProgressCb, temperature: float,
                    on_delta=None) -> str:
    """One provider call with the same streaming→non-streaming fallback the single-pass
    path uses, so a dropped stream still yields output."""
    try:
        return await ai_service.analyze_with_provider(
            full_prompt, system_prompt=system_prompt, provider=provider,
            on_progress=on_progress, on_delta=on_delta, temperature=temperature)
    except Exception:
        if on_delta is None:
            raise
        logger.exception("streamed LLM call failed — retrying without streaming")
        return await ai_service.analyze_with_provider(
            full_prompt, system_prompt=system_prompt, provider=provider,
            on_progress=on_progress, temperature=temperature)


# ── Per-slide QA ("test the presentation") ───────────────────────────────────────────────────
# Conservative, high-signal checks only, so a good slide is never churned. Runs per slide so a
# problem is repaired IN PLACE — a whole-document rewrite would throw away the other slides'
# compositions, which is the whole point of building them one at a time.
_SLIDE_PLOTLY_RE = re.compile(
    r'<script\b[^>]*\bclass=["\'][^"\']*\bplotly-spec\b[^"\']*["\'][^>]*>(.*?)</script>',
    re.IGNORECASE | re.DOTALL)
_STYLE_SCRIPT_STRIP_RE = re.compile(r'<(style|script)\b[^>]*>.*?</\1>', re.IGNORECASE | re.DOTALL)
_TAG_STRIP_RE = re.compile(r'<[^>]+>')
_BORDER_SIDE_RE = re.compile(r"border-(?:left|top|right|bottom)\s*:\s*\d", re.IGNORECASE)
_TEXT_OPACITY_RE = re.compile(r"opacity\s*:\s*0?\.\d", re.IGNORECASE)
_NESTED_CARD_RE = re.compile(r'class=["\'][^"\']*\bcard\b[^"\']*["\'][^>]*>(?:(?!</div>).)*?'
                             r'class=["\'][^"\']*\bcard\b', re.IGNORECASE | re.DOTALL)
_SLIDE_MAX_WORDS = 190


def _check_slide(section: str, *, with_photos: bool, wants_image: bool) -> List[str]:
    """Return this slide's problems (empty = clean). Mirrors deck_validation's philosophy but
    scoped to ONE slide, and adds the PowerPoint export-safety rules."""
    import json as _json
    probs: List[str] = []
    if not section or not _SECTION_RE.search(section):
        return ["The slide is not a single <section class=\"slide\">…</section>."]

    for i, body in enumerate(_SLIDE_PLOTLY_RE.findall(section), 1):
        raw = (body or "").strip()
        if not raw:
            probs.append(f"Chart spec #{i} is empty.")
            continue
        try:
            spec = _json.loads(raw)
        except Exception as e:
            probs.append(f"Chart spec #{i} is not valid JSON ({e}). Emit strictly valid JSON.")
            continue
        if not isinstance(spec, dict) or "data" not in spec or "layout" not in spec:
            probs.append(f'Chart spec #{i} is missing "data" and/or "layout".')

    visible = _TAG_STRIP_RE.sub(" ", _STYLE_SCRIPT_STRIP_RE.sub(" ", section))
    if '{"data"' in visible or '"layout":' in visible:
        probs.append("Chart JSON is visible on the slide — it must live ONLY inside the hidden "
                     "<script class=\"plotly-spec\"> element.")
    words = len(visible.split())
    if words > _SLIDE_MAX_WORDS:
        probs.append(f"The slide is overcrowded ({words} words) — cut to the essentials and let the "
                     "visual carry the detail.")

    # Export-safety (PowerPoint): these break the .pptx render.
    if "linear-gradient" in section and "photo-overlay" not in section and "ai-img" not in section:
        probs.append("Gradient background/panel found — export-safety forbids gradients (a scrim over "
                     "a cover photo is the only exception). Use a solid tint panel.")
    if _BORDER_SIDE_RE.search(section):
        probs.append("Single-side border/accent stripe found — use a full 4-side border or a solid "
                     "tint panel instead (export-safety).")
    if _TEXT_OPACITY_RE.search(section):
        probs.append("Translucent/faded element found (opacity < 1) — use a lighter palette colour "
                     "instead (export-safety).")
    if _NESTED_CARD_RE.search(section):
        probs.append("Nested cards (a .card inside a .card) — group with whitespace/headings instead.")

    # Images must match the plan's budget.
    has_img = "ai-img" in section
    if has_img and not with_photos:
        probs.append("This deck has photos disabled — remove the <img class=\"ai-img\"> placeholder.")
    elif has_img and not wants_image:
        probs.append("This slide has no planned image — remove the <img class=\"ai-img\"> placeholder.")
    return probs


_SLIDE_REPAIR_PROMPT = """The slide below has specific problems. Fix ONLY these problems and re-emit the
COMPLETE slide. Preserve its composition, content, wording and numbers exactly as they are.

PROBLEMS TO FIX:
{problems}

Output EXACTLY ONE <section class="slide {archetype}"> ... </section> and NOTHING ELSE — no markdown
fences, no <style> tag, no commentary. Keep using only the shared stylesheet's classes/tokens.

=== CURRENT SLIDE ===
{section}"""


_HEX_RE = re.compile(r"#[0-9A-Fa-f]{6}")


def _brand_hex(brand: Optional[str], which: str) -> str:
    """Pull the resolved accent hexes out of the brand directive the route built (it states
    '<hex> as --accent and <hex> as --accent-2'). Falls back to the TBS palette."""
    hexes = _HEX_RE.findall(brand or "")
    if which == "accent":
        return hexes[0] if hexes else TBS_PALETTE["accent"]
    return hexes[1] if len(hexes) > 1 else TBS_PALETTE["accent2"]


_NOTE_HEAD_RE = re.compile(r"^##\s*Slide\s*(\d+)\s*$", re.IGNORECASE | re.MULTILINE)


def _split_notes(raw: str, total: int) -> List[str]:
    """Split the batched speaker-notes reply into one entry per slide ('## Slide N' headings).
    Always returns exactly `total` entries so notes line up with slides positionally."""
    out = [""] * total
    if not raw:
        return out
    marks = list(_NOTE_HEAD_RE.finditer(raw))
    for i, m in enumerate(marks):
        end = marks[i + 1].start() if i + 1 < len(marks) else len(raw)
        try:
            idx = int(m.group(1)) - 1
        except ValueError:
            continue
        if 0 <= idx < total:
            out[idx] = raw[m.end():end].strip()
    return out


def _clean_css(text: str) -> str:
    """Strip fences / stray <style> wrappers from the stylesheet stage's reply."""
    t = (text or "").strip()
    if t.startswith("```"):
        t = t.split("\n", 1)[1] if "\n" in t else t
        if t.endswith("```"):
            t = t[: t.rfind("```")]
    t = re.sub(r"</?style[^>]*>", "", t, flags=re.IGNORECASE)
    return t.strip()


_SECTION_RE = re.compile(r"<section\b.*</section>", re.IGNORECASE | re.DOTALL)


def _clean_slide(text: str) -> str:
    """Extract exactly the <section>…</section> a per-slide call should have returned."""
    t = (text or "").strip()
    if t.startswith("```"):
        t = t.split("\n", 1)[1] if "\n" in t else t
        if t.endswith("```"):
            t = t[: t.rfind("```")]
    m = _SECTION_RE.search(t)
    return m.group(0).strip() if m else ""


def _fallback_slide(md: str, archetype: str, n: int, total: int) -> str:
    """A minimal but valid slide built from the markdown, used when ONE slide's HTML call fails —
    the deck still ships rather than the whole generation dying on a single page."""
    import html as _h
    lines = [l.strip() for l in (md or "").splitlines() if l.strip()][:12]
    body = "<br>".join(_h.escape(l) for l in lines) or "&nbsp;"
    return (f'<section class="slide {archetype}">\n'
            f'  <div style="display:flex;flex-direction:column;justify-content:center;height:1080px;padding:84px 108px">'
            f'{body}</div>\n'
            f'  <div class="pageno">{n:02d} / {total:02d}</div>\n</section>')


def _assemble_deck(shared_css: str, slides_html: List[str]) -> str:
    """Wrap the independently-generated slides in ONE document with the shared stylesheet, so the
    existing render/validate/theme path downstream sees a normal self-contained deck."""
    body = "\n".join(s for s in slides_html if s and s.strip())
    return ('<!DOCTYPE html>\n<html lang="en">\n<head>\n<meta charset="utf-8">\n<title>Report</title>\n'
            "<style>\n" + shared_css + "\n</style>\n</head>\n<body>\n" + body + "\n</body>\n</html>")


# How many per-slide calls may be in flight at once (DashScope rate limits + the shared event loop).
_SLIDE_CONCURRENCY = 3


async def _generate_per_slide(data_brief: str, *, brand: Optional[str], structure: Optional[str],
                              seed: Optional[str], creativity: str,
                              planner: str, insights_provider: str, html_provider: str,
                              temperature: float, on_progress: ProgressCb, on_delta,
                              variant_seed: Optional[str] = None, with_photos: bool = True,
                              style: Optional[str] = "tbs"):
    """Per-slide pipeline: guidelines → plan → shared stylesheet → per-slide md → per-slide HTML →
    assemble → speaker notes.

    Each slide gets its own call (its own full token budget + the model's full attention), which is what
    fixes composition — one call writing 15 slides rations both. Two things keep the independently-built
    slides coherent: the DESIGN GUIDELINES (a per-deck brief with palette semantics, export-safety and
    layout principles) and the SHARED STYLESHEET (slides may only use ITS classes/tokens). Crucially the
    PLAN specifies each slide's LAYOUT composition, so the HTML stage executes a designed composition
    rather than improvising geometry.

    Returns (html_document, artifacts_dict). Raises on plan failure so the caller can fall back to
    single-pass."""
    import json
    directive = _structure_directive(creativity, structure or DEFAULT_STRUCTURE)
    art = _variant_directive(variant_seed or seed or "")
    style_dir = _style_directive(style, seed)

    # ── Stage 0: design guidelines — the brief every later stage is built from ──
    if on_progress:
        await on_progress("Setting the design direction…")
    guidelines = (await _call_llm(
        "\n\n".join([
            _GUIDELINES_PROMPT
            .replace("{canvas}", f"{SLIDE_W_PX}x{SLIDE_H_PX}")
            .replace("{accent}", _brand_hex(brand, "accent"))
            .replace("{accent2}", _brand_hex(brand, "accent2"))
            .replace("{data}", data_brief),
            style_dir, art,
        ]),
        system_prompt=_GUIDELINES_SYSTEM, provider=planner, on_progress=None, temperature=0.5)).strip()

    # ── Stage 1: plan (designs each slide's LAYOUT, not just its content) ──
    if on_progress:
        await on_progress("Planning slides…")
    max_images = MAX_DECK_IMAGES if with_photos else 0
    plan_raw = await _call_llm(
        _PLAN_PROMPT.replace("{structure_directive}", directive)
        .replace("{guidelines}", guidelines)
        .replace("{max_images}", str(max_images))
        .replace("{data}", data_brief),
        system_prompt=_PLAN_SYSTEM, provider=planner, on_progress=None, temperature=0.5)
    plan = _extract_json(plan_raw)
    slides = (plan or {}).get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("planner did not return a usable slide plan")
    plan_json = json.dumps(plan, ensure_ascii=False, indent=1)
    total = len(slides)

    # ── Stage 2: shared stylesheet (the consistency backbone) ──
    if on_progress:
        await on_progress("Designing the system…")
    css_prompt = "\n\n".join([_STYLESHEET_PROMPT, DESIGN_SYSTEM, THEME_PRESETS, style_dir, art,
                              (brand or ""), "=== DESIGN GUIDELINES ===\n" + guidelines])
    shared_css = _clean_css(await _call_llm(css_prompt, system_prompt=_STYLESHEET_SYSTEM,
                                            provider=html_provider, on_progress=None, temperature=0.5))

    sem = asyncio.Semaphore(_SLIDE_CONCURRENCY)

    # ── Stage 3: per-slide markdown ──
    written = {"n": 0}

    async def _md(slide: dict) -> str:
        async with sem:
            try:
                out = await _call_llm(
                    _SLIDE_MD_PROMPT.replace("{slide_json}", json.dumps(slide, ensure_ascii=False))
                    .replace("{plan}", plan_json).replace("{data}", data_brief),
                    system_prompt=_INSIGHTS_SYSTEM, provider=insights_provider,
                    on_progress=None, temperature=0.6)
            except Exception:
                logger.exception("slide md failed (n=%s) — using the plan entry as copy", slide.get("n"))
                out = f"HEADLINE: {slide.get('title','')}\nTAKEAWAY: {slide.get('purpose','')}"
            written["n"] += 1
            if on_progress:
                await on_progress(f"Writing slide {written['n']}/{total}…")
            return (out or "").strip()

    mds = list(await asyncio.gather(*[_md(s) for s in slides]))

    # ── Stage 4: per-slide HTML — EXECUTE the planned composition ──
    built = {"n": 0}

    def _photo_rule_for(slide: dict) -> str:
        """Photos come from the PLAN's image budget, not the designer's whim — a data-led deck gets ~2."""
        if not with_photos:
            return ('PHOTO: do NOT use any <img class="ai-img"> — build this slide from colour fields, '
                    'typography, charts and icons only.')
        img = slide.get("image") or {}
        if img.get("needed") and img.get("prompt"):
            return (f'PHOTO: this slide HAS a planned image. Add <img class="ai-img" '
                    f'data-prompt="{img.get("prompt")}"> (no src — the system fills it), '
                    f'{img.get("sizing") or "object-cover"}, positioned per the LAYOUT above; put a scrim '
                    f'under any text sitting on it.')
        return ('PHOTO: this slide has NO planned image — do NOT add one. Never render a chart or table '
                'as an image.')

    async def _html(slide: dict, md: str, idx: int) -> str:
        archetype = slide.get("archetype") or "layout-split"
        n = idx + 1
        blocks = slide.get("blocks") or []
        blocks_txt = "\n".join(
            f"- [{b.get('type','text')}] @ {b.get('placement','')}: {b.get('content','')}"
            + (f"  (brief: {b.get('creative_brief')})" if b.get("creative_brief") else "")
            for b in blocks if isinstance(b, dict)
        ) or "- (none specified — compose from the slide copy below)"
        async with sem:
            try:
                prompt_txt = (_SLIDE_HTML_PROMPT
                              .replace("{layout}", slide.get("layout") or
                                       f"Follow the {archetype} archetype from the design system.")
                              .replace("{blocks}", blocks_txt)
                              .replace("{md}", md)
                              .replace("{archetype}", archetype)
                              .replace("{index_label}", f"{n:02d} / {total:02d}")
                              .replace("{n}", str(n)).replace("{total}", str(total))
                              .replace("{icons}", _SLIDE_ICON_NAMES)
                              .replace("{photo_rule}", _photo_rule_for(slide))
                              .replace("{guidelines}", guidelines)
                              .replace("{art_direction}", art)
                              .replace("{shared_css}", shared_css))
                raw = await _call_llm(prompt_txt, system_prompt=_SLIDE_HTML_SYSTEM,
                                      provider=html_provider, on_progress=None,
                                      temperature=temperature, on_delta=on_delta)
                out = _clean_slide(raw) or _fallback_slide(md, archetype, n, total)
            except Exception:
                logger.exception("slide html failed (n=%s) — using a fallback slide", n)
                out = _fallback_slide(md, archetype, n, total)
            built["n"] += 1
            if on_progress:
                await on_progress(f"Building slide {built['n']}/{total}…")
            return out

    htmls = list(await asyncio.gather(*[_html(s, m, i) for i, (s, m) in enumerate(zip(slides, mds))]))

    # ── Stage 4b: QA each slide and repair ONLY the broken ones, in place ──
    if on_progress:
        await on_progress("Checking slides…")

    async def _qa(slide: dict, section: str, idx: int) -> str:
        n = idx + 1
        wants_image = bool((slide.get("image") or {}).get("needed"))
        probs = _check_slide(section, with_photos=with_photos, wants_image=wants_image)
        if not probs:
            return section
        logger.info("slide %d/%d failed QA: %s", n, total, "; ".join(probs))
        archetype = slide.get("archetype") or "layout-split"
        async with sem:
            try:
                fixed_raw = await _call_llm(
                    _SLIDE_REPAIR_PROMPT
                    .replace("{problems}", "\n".join(f"- {p}" for p in probs))
                    .replace("{archetype}", archetype)
                    .replace("{section}", section),
                    system_prompt=_SLIDE_HTML_SYSTEM, provider=html_provider,
                    on_progress=None, temperature=0.4, on_delta=on_delta)
                fixed = _clean_slide(fixed_raw)
            except Exception:
                logger.exception("slide %d repair call failed — keeping the original", n)
                return section
        # Only accept a repair that is itself clean (or at least no worse) — never ship a
        # regression from the fixer.
        if fixed and len(_check_slide(fixed, with_photos=with_photos, wants_image=wants_image)) < len(probs):
            if on_progress:
                await on_progress(f"Fixed slide {n}/{total}…")
            return fixed
        return section

    htmls = list(await asyncio.gather(*[_qa(s, h, i) for i, (s, h) in enumerate(zip(slides, htmls))]))
    doc = _assemble_deck(shared_css, htmls)

    # ── Stage 5: speaker notes (one batched call, after the slides exist) ──
    if on_progress:
        await on_progress("Writing speaker notes…")
    notes: List[str] = []
    try:
        copy_txt = "\n\n".join(f"## Slide {i + 1}\n{m}" for i, m in enumerate(mds))
        notes_raw = await _call_llm(
            _NOTES_PROMPT.replace("{plan}", plan_json).replace("{copy}", copy_txt).replace("{n}", "N"),
            system_prompt=_NOTES_SYSTEM, provider=insights_provider, on_progress=None, temperature=0.6)
        notes = _split_notes(notes_raw, total)
    except Exception:
        logger.exception("speaker notes failed — shipping the deck without them")

    return doc, {"guidelines_md": guidelines, "slide_plan": plan,
                 "slides_md": mds, "slides_html": htmls, "slides_notes": notes}


async def generate_deck_html(data_brief: str, *, prompt: Optional[str] = None,
                             brand: Optional[str] = None, structure: Optional[str] = None,
                             provider: str = "qwen3.7-max", on_progress: ProgressCb = None,
                             image_cache: Optional[Dict[str, "asyncio.Task"]] = None,
                             seed: Optional[str] = None, creativity: str = "balanced",
                             pipeline: str = "single", models: Optional[Dict[str, str]] = None,
                             style: Optional[str] = "tbs",
                             artifacts: Optional[Dict] = None) -> str:
    """Ask the chosen LLM provider to design the deck and return self-contained HTML.

    `artifacts`, when given, is filled with the per-slide {slides_md, slides_html} produced by the
    per-slide pipeline so the caller can persist one artifact per page.

    Runs one cheap (no-browser) validation pass; if it finds structural, Plotly-spec,
    or ungrounded-number problems, makes a single targeted repair call. Always returns
    renderable HTML best-effort — a failed/worse repair is discarded, never blocks.

    If image_cache is provided, the deck is streamed and image generation is kicked off
    for each photo placeholder as it appears (concurrent with writing); resolve_ai_images
    later awaits those tasks. Falls back transparently to non-streaming if streaming fails."""
    from services.deck_validation import validate_deck_html
    from services.image_service import images_enabled

    # More creative freedom → a bit more sampling variety.
    temperature = {"structured": 0.7, "balanced": 0.85, "creative": 1.0}.get(creativity, 0.85)
    # Stream + prewarm images only when a cache is supplied and images are enabled.
    on_delta = (_make_image_prewarmer(image_cache)
                if image_cache is not None and images_enabled() else None)
    # Per-layer model overrides (3-layer pipeline); each defaults to `provider`.
    models = models or {}
    html_provider = models.get("html") or provider
    # Art-direction seed is per-GENERATION (+ the html model), so two models — or two runs — of the
    # SAME site produce different compositions. Theme/fonts/brand still pin to `seed` (the domain).
    import uuid
    variant_seed = f"{html_provider}:{uuid.uuid4().hex}"
    photos_on = image_cache is not None and images_enabled()  # AI-photo toggle for this deck

    async def _single_pass() -> str:
        # prompt is normally resolved by the route (from the chosen prompt id); fall back
        # to the built-in default if none was passed.
        full_prompt = build_prompt(data_brief, prompt=prompt, brand=brand, structure=structure,
                                   seed=seed, creativity=creativity, variant_seed=variant_seed,
                                   style=style)
        if on_progress:
            await on_progress("Writing slides…")
        return await _call_llm(full_prompt, system_prompt=_DECK_SYSTEM_PROMPT, provider=provider,
                               on_progress=on_progress, temperature=temperature, on_delta=on_delta)

    if pipeline == "layered":
        try:
            raw, produced = await _generate_per_slide(
                data_brief, brand=brand, structure=structure, seed=seed, creativity=creativity,
                planner=models.get("planner") or provider,
                insights_provider=models.get("insights") or provider,
                html_provider=html_provider, temperature=temperature,
                on_progress=on_progress, on_delta=on_delta, variant_seed=variant_seed,
                with_photos=photos_on, style=style)
            # Hand the per-page artifacts back to the caller so they can be stored on the Document
            # (design guidelines, the slide plan, and one .md/.html/notes per slide).
            if artifacts is not None:
                artifacts.update(produced)
        except Exception:
            logger.exception("per-slide pipeline failed — falling back to single-pass")
            raw = await _single_pass()
    else:
        raw = await _single_pass()
    html = _clean_html(raw)

    result = validate_deck_html(html, data_brief, creativity=creativity)
    if result.ok:
        return _strip_leaked_specs(html)

    # The per-slide pipeline already QA'd and repaired each slide individually. A whole-document
    # repair here would rewrite every slide at once — throwing away the per-slide compositions that
    # are the entire point of building them one at a time. So unless the assembled doc is
    # STRUCTURALLY broken, log the remaining notes and ship it.
    if pipeline == "layered" and not result.structural:
        logger.info("layered deck: per-slide QA passed; skipping whole-document repair "
                    "(plotly=%d, ungrounded_numbers=%d, design=%d)",
                    len(result.plotly), len(result.ungrounded_numbers), len(result.design))
        return _strip_leaked_specs(html)

    logger.info(
        "deck validation found issues (structural=%d, plotly=%d, ungrounded_numbers=%d) — attempting repair",
        len(result.structural), len(result.plotly), len(result.ungrounded_numbers),
    )
    if on_progress:
        await on_progress("Checking & fixing the deck…")
    try:
        repair_raw = await ai_service.analyze_with_provider(
            _build_repair_prompt(html, result.repair_instructions()),
            system_prompt=_DECK_SYSTEM_PROMPT,
            provider=html_provider,
            on_progress=on_progress,
        )
        repaired = _clean_html(repair_raw)
        after = validate_deck_html(repaired, data_brief)
        # Accept the repair only if it is still a valid full deck (no structural errors),
        # so a bad repair can never leave us worse off than the original.
        if not after.structural:
            html = repaired
            logger.info(
                "post-repair validation (structural=%d, plotly=%d, ungrounded_numbers=%d)",
                len(after.structural), len(after.plotly), len(after.ungrounded_numbers),
            )
        else:
            logger.warning("repair produced a structurally invalid deck — keeping original HTML")
    except Exception:
        logger.exception("deck repair call failed — keeping original HTML")

    return _strip_leaked_specs(html)


# ---- Rendering (headless Chromium via Playwright) --------------------------
_PLOTLY_SRC = None


def _plotly_source() -> str:
    """Load the bundled Plotly library once (backend/assets/plotly.min.js)."""
    global _PLOTLY_SRC
    if _PLOTLY_SRC is None:
        f = Path(__file__).resolve().parent.parent / "assets" / "plotly.min.js"
        try:
            _PLOTLY_SRC = f.read_text(encoding="utf-8")
        except Exception:
            _PLOTLY_SRC = ""
    return _PLOTLY_SRC


_PLOTLY_SPEC_RE = re.compile(
    r'(<script\b[^>]*\bclass=["\']plotly-spec["\'][^>]*>)(.*?)(</script>)',
    re.IGNORECASE | re.DOTALL,
)
# How big the largest bubble should render (px diameter) and how many to label.
_BUBBLE_MAX_PX = 55
_BUBBLE_LABELS = 8


def _is_keyword_bubble(trace: dict, layout: dict) -> bool:
    """Identify the keyword position-vs-impressions scatter so we can enforce its
    sizing/labels. Prefer an explicit meta flag the prompt asks for; else fall back to
    the axis titles the model reliably emits (x≈'position', y≈'impressions')."""
    if (trace.get("meta") or {}).get("chart") == "keyword-bubble":
        return True
    if (layout.get("meta") or {}).get("chart") == "keyword-bubble":
        return True
    if str(trace.get("type", "scatter")).lower() not in ("scatter", "scattergl", ""):
        return False

    def _title(ax):
        t = (layout.get(ax) or {}).get("title")
        return (t.get("text") if isinstance(t, dict) else t) or ""

    xt, yt = _title("xaxis").lower(), _title("yaxis").lower()
    return "position" in xt and "impress" in yt


def _enforce_keyword_bubble(html: str) -> str:
    """Deterministically fix the keyword bubble chart regardless of what the LLM emitted:
    scale every bubble's area by impressions and label the top-N by impressions beside their
    point. The model only needs to supply the data (x positions, y impressions, and a query
    per point via customdata/text/hovertext); the math + labelling happen here so the chart
    is reliably big and labelled even when the model ignores those instructions."""
    if "plotly-spec" not in html:
        return html
    import json

    def _floats(seq):
        out = []
        for v in seq:
            try:
                out.append(float(v))
            except (TypeError, ValueError):
                out.append(0.0)
        return out

    def _fix(m):
        try:
            spec = json.loads(m.group(2))
        except Exception:
            return m.group(0)
        data = spec.get("data")
        layout = spec.get("layout") or {}
        if not isinstance(data, list):
            return m.group(0)

        changed = False
        for trace in data:
            if not isinstance(trace, dict) or "y" not in trace or not isinstance(trace.get("y"), list):
                continue
            if not _is_keyword_bubble(trace, layout):
                continue
            ys = _floats(trace["y"])
            if not ys or max(ys) <= 0:
                continue
            maxy = max(ys)
            marker = trace.get("marker") if isinstance(trace.get("marker"), dict) else {}
            marker["size"] = ys
            marker["sizemode"] = "area"
            marker["sizeref"] = (2.0 * maxy) / (_BUBBLE_MAX_PX ** 2)
            marker["sizemin"] = 4
            marker.setdefault("opacity", 0.7)
            trace["marker"] = marker

            # Per-point query labels — pull from whatever the model provided.
            labels_src = None
            for key in ("customdata", "text", "hovertext"):
                val = trace.get(key)
                if isinstance(val, list) and len(val) == len(ys):
                    labels_src = [("" if v is None else str(v)) for v in val]
                    break
            if labels_src:
                keep = set(sorted(range(len(ys)), key=lambda i: ys[i], reverse=True)[:_BUBBLE_LABELS])
                trace["text"] = [labels_src[i] if i in keep else "" for i in range(len(ys))]
                trace["mode"] = "markers+text"
                # Place each label on the INWARD side so it can't run off the slide edge:
                # the x-axis is reversed (best ranks on the right), so points with a small x
                # sit on the right and get a left-side label; the rest get a right-side label.
                xs = _floats(trace.get("x") or list(range(len(ys))))
                xmid = (min(xs) + max(xs)) / 2 if xs else 0
                trace["textposition"] = [
                    ("middle left" if xs[i] <= xmid else "middle right") for i in range(len(ys))
                ]
                tf = trace.get("textfont") if isinstance(trace.get("textfont"), dict) else {}
                tf.setdefault("size", 11)
                trace["textfont"] = tf
            else:
                trace["mode"] = "markers"

            # Headroom on BOTH sides so inward labels (left or right) never clip.
            margin = layout.get("margin") if isinstance(layout.get("margin"), dict) else {}
            margin["r"] = max(int(margin.get("r", 0) or 0), 130)
            margin["l"] = max(int(margin.get("l", 0) or 0), 120)
            layout["margin"] = margin
            spec["layout"] = layout
            changed = True

        if not changed:
            return m.group(0)
        return m.group(1) + json.dumps(spec) + m.group(3)

    return _PLOTLY_SPEC_RE.sub(_fix, html)


_ACCENT_VAR_RE = re.compile(r'(--accent\s*:\s*)([^;}\n]+)([;}])', re.IGNORECASE)
_ACCENT2_VAR_RE = re.compile(r'(--accent-2\s*:\s*)([^;}\n]+)([;}])', re.IGNORECASE)


def _apply_theme(html: str, accent: str, accent2: str, muted: str = "#9AA0A6") -> str:
    """Force the deck onto the site's brand colour, deterministically — so themeing never
    depends on the model obeying the prompt. Rewrites the --accent / --accent-2 CSS tokens
    (covering every var(--accent) use in the slide design) and recolours every Plotly series
    to accent / accent-2 (combo bars=accent, line=accent-2; pie = accent/accent-2/muted;
    choropleth = a light→accent ramp). Only colour keys are touched; data is never changed."""
    if not accent:
        return html
    html = _ACCENT_VAR_RE.sub(lambda m: m.group(1) + accent + m.group(3), html, count=1)
    html = _ACCENT2_VAR_RE.sub(lambda m: m.group(1) + accent2 + m.group(3), html, count=1)
    if "plotly-spec" not in html:
        return html
    import json

    palette = [accent, accent2]

    def _recolor(m):
        try:
            spec = json.loads(m.group(2))
        except Exception:
            return m.group(0)
        data = spec.get("data")
        if not isinstance(data, list):
            return m.group(0)
        ci = 0
        for trace in data:
            if not isinstance(trace, dict):
                continue
            ttype = str(trace.get("type", "")).lower()
            if ttype == "pie":
                mk = trace.get("marker") if isinstance(trace.get("marker"), dict) else {}
                vals = trace.get("values") or []
                base = [accent, accent2, muted]
                mk["colors"] = [base[i % 3] for i in range(len(vals) or 3)]
                trace["marker"] = mk
            elif ttype == "choropleth":
                trace["colorscale"] = [[0, "#F1EDE6"], [1, accent]]
                trace.setdefault("reversescale", False)
            else:
                color = palette[ci % len(palette)]
                ci += 1
                mk = trace.get("marker") if isinstance(trace.get("marker"), dict) else {}
                mk["color"] = color
                trace["marker"] = mk
                if isinstance(trace.get("line"), dict):
                    trace["line"]["color"] = color
                elif str(trace.get("mode", "")).find("lines") != -1:
                    trace["line"] = {"color": color}
        return m.group(1) + json.dumps(spec) + m.group(3)

    return _PLOTLY_SPEC_RE.sub(_recolor, html)


def _polish_plotly_specs(html: str) -> str:
    """Apply one cohesive 'house style' to EVERY chart so the deck looks like a designed
    template, not N improvised charts — done deterministically (setdefault only, so explicit
    model choices win) instead of trusting the LLM to style each chart. Fixes the recurring
    rough edges: transparent backgrounds, faint gridlines, tidy margins, rounded bars,
    automargin so axis labels stop overlapping, and compact choropleth colorbars."""
    if "plotly-spec" not in html:
        return html
    import json

    def _polish(m):
        try:
            spec = json.loads(m.group(2))
        except Exception:
            return m.group(0)
        data = spec.get("data")
        layout = spec.get("layout") if isinstance(spec.get("layout"), dict) else {}
        if not isinstance(data, list):
            return m.group(0)

        layout.setdefault("paper_bgcolor", "rgba(0,0,0,0)")
        layout.setdefault("plot_bgcolor", "rgba(0,0,0,0)")
        layout.setdefault("hovermode", False)
        margin = layout.get("margin") if isinstance(layout.get("margin"), dict) else {}
        for k, v in (("t", 16), ("r", 16), ("b", 44), ("l", 52)):
            margin.setdefault(k, v)
        layout["margin"] = margin
        font = layout.get("font") if isinstance(layout.get("font"), dict) else {}
        font.setdefault("size", 13)
        layout["font"] = font
        # Every axis (xaxis, yaxis, xaxis2, …) gets automargin + clean grid/zeroline.
        for key, ax in list(layout.items()):
            if (key.startswith("xaxis") or key.startswith("yaxis")) and isinstance(ax, dict):
                ax.setdefault("automargin", True)
                ax.setdefault("zeroline", False)
                ax.setdefault("gridcolor", "rgba(0,0,0,0.06)")
                ax.setdefault("linecolor", "rgba(0,0,0,0.15)")
                ax.setdefault("ticks", "")

        for tr in data:
            if not isinstance(tr, dict):
                continue
            ttype = str(tr.get("type", "")).lower()
            if ttype == "bar":
                mk = tr.get("marker") if isinstance(tr.get("marker"), dict) else {}
                mk.setdefault("cornerradius", 5)
                tr["marker"] = mk
            elif ttype in ("scatter", "scattergl", "") and isinstance(tr.get("line"), dict):
                tr["line"].setdefault("width", 2.5)
            elif ttype == "choropleth":
                cb = tr.get("colorbar") if isinstance(tr.get("colorbar"), dict) else {}
                cb.setdefault("thickness", 12)
                cb.setdefault("len", 0.6)
                cb.setdefault("outlinewidth", 0)
                cb.setdefault("x", 1.02)
                tr["colorbar"] = cb

        spec["layout"] = layout
        return m.group(1) + json.dumps(spec) + m.group(3)

    return _PLOTLY_SPEC_RE.sub(_polish, html)


# Deterministic full-height CSS, injected last with !important so it wins over the model's
# inline styles. Forces every slide to a full-height flex column and vertically centers its
# content, so an under-filled slide rebalances its whitespace (no one-sided blank band)
# instead of clustering at the top. Charts that sit in a growing flex region are stretched to
# fill — but fixed-aspect charts (pie/choropleth) keep their own height so they don't distort.
_FILL_CSS = """<style>/* deterministic-fill */
.slide{display:flex !important;flex-direction:column !important;height:1080px !important;
  justify-content:center !important;overflow:hidden !important;position:relative !important;}
/* pin the page index to the bottom-right so centering the content can't displace it */
.slide .pageno{position:absolute !important;right:64px;bottom:30px;margin:0 !important;}
/* let a chart that is the slide's direct hero grow to fill the column */
.slide > .js-plotly-plot{flex:1 1 auto;min-height:320px;}
</style>"""


def _enforce_fill(html: str) -> str:
    """Append the full-height fill stylesheet just before </head> (or </body> as a fallback)
    so it loads after the deck's own <style> and wins on cascade order + !important."""
    if "slide" not in html:
        return html
    if "</head>" in html:
        return html.replace("</head>", _FILL_CSS + "</head>", 1)
    if "</body>" in html:
        return html.replace("</body>", _FILL_CSS + "</body>", 1)
    return html + _FILL_CSS


# Render-time DOM cleanup: in the parsed page (where escaping is resolved and "visible" is
# unambiguous) remove any TEXT node the model printed that is a raw Plotly chart spec — a JSON
# object carrying "data" + ("layout" | paper_bgcolor). This is the authoritative guard against
# leaked specs: it catches HTML-escaped, key-reordered and oddly-wrapped copies the pre-render
# regex (_strip_leaked_specs) can't. It SKIPS <script>/<style> so the real plotly-spec scripts
# Plotly reads survive, and never removes an element with an id (e.g. a chart container).
_LEAK_CLEANUP_JS = (
    "<script>(function(){try{"
    "function isSpec(t){t=(t||'').trim();"
    "return t.length>40&&t.charAt(0)==='{'&&t.indexOf('\"data\"')!==-1&&"
    "(t.indexOf('\"layout\"')!==-1||t.indexOf('paper_bgcolor')!==-1);}"
    "var w=document.createTreeWalker(document.body,NodeFilter.SHOW_TEXT,null),n,kill=[];"
    "while(n=w.nextNode()){var p=n.parentNode;if(!p)continue;"
    "var tag=p.nodeName;if(tag==='SCRIPT'||tag==='STYLE')continue;"
    "if(isSpec(n.nodeValue))kill.push(n);}"
    "kill.forEach(function(n){var p=n.parentNode;n.nodeValue='';"
    "if(p&&p!==document.body&&!p.id&&!p.querySelector('*')&&!(p.textContent||'').trim()){"
    "try{p.parentNode.removeChild(p);}catch(e){}}});"
    "}catch(e){}})();</script>"
)


def _inject_leak_cleanup(html: str) -> str:
    """Inject the DOM leak-cleanup script before </body> so it runs on every deck (independent
    of whether charts are present)."""
    if "</body>" in html:
        return html.replace("</body>", _LEAK_CLEANUP_JS + "</body>", 1)
    return html + _LEAK_CLEANUP_JS


def _prepare_html_for_render(html: str) -> str:
    """Inline the bundled Plotly so charts render with NO internet access — the VPS
    may not reach cdn.plot.ly. Strips the CDN <script> and injects the local copy."""
    # Inline fonts first so typography is offline-safe even for text-only (chart-less) decks.
    html = _inline_fonts(html)
    # Deterministic full-height backstop — applies to every deck (incl. chart-less) so a
    # model that under-fills a slide still produces a balanced, full-canvas page.
    html = _enforce_fill(html)
    # Authoritative leaked-spec guard — runs on every deck, including the chart-less path.
    html = _inject_leak_cleanup(html)
    if not any(k in html for k in ("Plotly.newPlot", "plot.ly", "plotly-spec")):
        return html
    html = _polish_plotly_specs(html)
    html = _enforce_keyword_bubble(html)
    src = _plotly_source()
    if not src:
        return html  # no local bundle available — leave the CDN tag as-is
    html = re.sub(r'<script\b[^>]*\bsrc=["\']https?://cdn\.plot\.ly/[^"\']*["\'][^>]*>\s*</script>',
                  '', html, flags=re.IGNORECASE)
    inline = "<script>/* bundled-plotly */\n" + src + "\n</script>"
    if "</head>" in html:
        html = html.replace("</head>", inline + "</head>", 1)
    else:
        html = inline + html
    # Render every <script type="application/json" class="plotly-spec"> ourselves — the
    # AI emits JSON specs reliably but inconsistently wires up executable <script> calls.
    bootstrap = (
        "<script>(function(){if(!window.Plotly)return;"
        "document.querySelectorAll('script.plotly-spec').forEach(function(s){try{"
        "var spec=JSON.parse(s.textContent);"
        "var id=s.getAttribute('data-target');"
        "var el=(id&&document.getElementById(id))||s.previousElementSibling;"
        "if(el)Plotly.newPlot(el,spec.data||[],spec.layout||{},"
        "{staticPlot:true,displayModeBar:false,responsive:false});"
        "}catch(e){console.error('plotly-spec render failed',e);}});})();</script>"
    )
    if "</body>" in html:
        return html.replace("</body>", bootstrap + "</body>", 1)
    return html + bootstrap


_FONTS_CSS: Optional[str] = None
_GOOGLE_FONT_LINK_RE = re.compile(
    r'<link\b[^>]*fonts\.googleapis\.com[^>]*>', re.IGNORECASE)
_GOOGLE_FONT_IMPORT_RE = re.compile(
    r'@import\s+url\([^)]*fonts\.googleapis\.com[^)]*\)\s*;?', re.IGNORECASE)


def _fonts_css() -> str:
    """Load the bundled font CSS once and bind its {FONTS_DIR} placeholder to the absolute
    file URI of assets/fonts/ so the woff2 files resolve from the temp HTML file. Empty if
    the bundle hasn't been fetched (then the remote Google Font link is left in place)."""
    global _FONTS_CSS
    if _FONTS_CSS is None:
        assets = Path(__file__).resolve().parent.parent / "assets"
        f = assets / "fonts.css"
        try:
            uri = (assets / "fonts").as_uri()  # file:///C:/.../assets/fonts
            _FONTS_CSS = f.read_text(encoding="utf-8").replace("{FONTS_DIR}", uri)
        except Exception:
            _FONTS_CSS = ""
    return _FONTS_CSS


def _inline_fonts(html: str) -> str:
    """Inject the bundled @font-face CSS and drop the remote Google Font references so the
    deck renders with correct typography even on a host that can't reach fonts.googleapis.com.
    No-op if the bundle is missing — the remote <link> stays so online rendering still works."""
    css = _fonts_css()
    if not css:
        return html
    html = _GOOGLE_FONT_LINK_RE.sub("", html)
    html = _GOOGLE_FONT_IMPORT_RE.sub("", html)
    block = "<style>/* bundled-fonts */\n" + css + "\n</style>"
    if "<head>" in html:
        return html.replace("<head>", "<head>" + block, 1)
    if "</head>" in html:
        return html.replace("</head>", block + "</head>", 1)
    return block + html


_AI_IMG_RE = re.compile(r'<img\b[^>]*\bai-img\b[^>]*>', re.IGNORECASE)


async def resolve_ai_images(html: str, *, max_images: int = MAX_DECK_IMAGES,
                            on_progress: ProgressCb = None,
                            image_cache: Optional[Dict[str, "asyncio.Task"]] = None) -> str:
    """Replace <img class="ai-img" data-prompt="..."> placeholders with real
    Supermachine photos embedded as base64 (so the offline renderer can use them).
    Failed or over-cap placeholders are dropped so the deck still renders cleanly.

    If image_cache is given, photos whose generation was already kicked off during the
    streamed write are awaited from it (usually already done) instead of being generated
    from scratch here — that's what overlaps image generation with slide-writing."""
    if "ai-img" not in html:
        return html
    from services.image_service import generate_image, images_enabled

    tags = list(_AI_IMG_RE.finditer(html))
    if not tags:
        return html
    if not images_enabled():
        # no image key — strip placeholders so no empty <img> shows
        return _AI_IMG_RE.sub("", html)

    gen_tags = tags[:max_images]
    if on_progress:
        await on_progress(f"Finalizing images… ({len(gen_tags)})")
    sem = asyncio.Semaphore(_IMG_CONCURRENCY)

    async def _one(prompt: str):
        # Prefer the prewarmed task started while the deck was being written.
        if image_cache is not None and prompt in image_cache:
            try:
                return await image_cache[prompt]
            except Exception:
                return None
        async with sem:
            return await generate_image(prompt)

    images = await asyncio.gather(*[_one(_img_prompt(m.group(0))) for m in gen_tags])

    # Clean up any prewarmed tasks the final deck didn't use (e.g. a placeholder the
    # repair pass changed) so they don't leak as 'never retrieved' or keep running.
    if image_cache:
        used = {_img_prompt(m.group(0)) for m in gen_tags}
        for p, task in image_cache.items():
            if p not in used and not task.done():
                task.cancel()
        await asyncio.gather(*image_cache.values(), return_exceptions=True)

    out, last = [], 0
    for idx, m in enumerate(tags):
        out.append(html[last:m.start()])
        img = images[idx] if idx < len(gen_tags) else None
        if img:
            mime = "jpeg" if img[:2] == b"\xff\xd8" else "png"
            b64 = base64.b64encode(img).decode()
            out.append(m.group(0).replace("<img", f'<img src="data:image/{mime};base64,{b64}"', 1))
        # else: drop the placeholder (failed or beyond cap)
        last = m.end()
    out.append(html[last:])
    return "".join(out)


_AI_ICON_RE = re.compile(r'<i\b[^>]*\bai-icon\b[^>]*>\s*</i>', re.IGNORECASE)
_ICONS_SRC: Optional[Dict[str, str]] = None


def _icons_source() -> Dict[str, str]:
    """Load the bundled inline-SVG icon set once (backend/assets/icons.json, Lucide/ISC)."""
    global _ICONS_SRC
    if _ICONS_SRC is None:
        import json
        f = Path(__file__).resolve().parent.parent / "assets" / "icons.json"
        try:
            _ICONS_SRC = json.loads(f.read_text(encoding="utf-8"))
        except Exception:
            _ICONS_SRC = {}
    return _ICONS_SRC


_ICON_COMMENT_RE = re.compile(r'<!--.*?-->', re.DOTALL)
_ICON_SIZE_RE = re.compile(r'\b(width|height)="[^"]*"', re.IGNORECASE)


def _prep_icon_svg(svg: str) -> str:
    """Make a vendored Lucide SVG inline-friendly: drop the license comment, size it to
    1em (so font-size controls it) and tag it .ai-icon. Stroke is already currentColor,
    so it inherits the surrounding text colour (set to --accent/--ink via CSS)."""
    svg = _ICON_COMMENT_RE.sub("", svg).strip()
    svg = _ICON_SIZE_RE.sub(lambda m: f'{m.group(1)}="1em"', svg)
    # ensure our class is present for sizing/vertical-align in the deck CSS
    if "ai-icon" not in svg:
        svg = svg.replace("<svg", '<svg class="ai-icon"', 1)
    return svg


def resolve_ai_icons(html: str) -> str:
    """Replace <i class="ai-icon" data-icon="NAME"></i> placeholders with inline SVG from
    the bundled set. Unknown icons are dropped (no broken glyphs). No network — instant,
    offline-safe — so this runs synchronously at the end of the deck pipeline."""
    if "ai-icon" not in html:
        return html
    icons = _icons_source()
    if not icons:
        return _AI_ICON_RE.sub("", html)

    def _sub(m: "re.Match") -> str:
        tag = m.group(0)
        name_m = re.search(r'data-icon=["\']([\w-]+)["\']', tag, re.IGNORECASE)
        svg = icons.get(name_m.group(1)) if name_m else None
        return _prep_icon_svg(svg) if svg else ""

    return _AI_ICON_RE.sub(_sub, html)


async def _render(html: str) -> Dict:
    """Render HTML once; return {'pdf': bytes, 'slides': [png bytes, ...]}.

    Loads via goto(temp file) rather than set_content: set_content injects HTML via
    document.write, which makes Chromium intermittently block the cross-site Plotly
    CDN script (leaving empty charts). A real navigation loads it reliably.
    """
    import os
    import tempfile
    from pathlib import Path
    from playwright.async_api import async_playwright

    tmp = tempfile.NamedTemporaryFile(mode="w", suffix=".html", delete=False, encoding="utf-8")
    try:
        tmp.write(_prepare_html_for_render(html))
        tmp.close()
        uses_plotly = "Plotly.newPlot" in html or "plot.ly" in html
        async with async_playwright() as p:
            browser = await p.chromium.launch(args=["--no-sandbox"])
            page = await browser.new_page(viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX},
                                          device_scale_factor=_RENDER_SCALE)
            await page.goto(Path(tmp.name).as_uri(), wait_until="networkidle")
            # Pin SCREEN media for the whole pipeline. The deck is designed for screen
            # (dark/saturated colour fields, full-bleed photos, 1080px flex columns); print
            # emulation would drop those, which is what made the old page.pdf() output blank.
            await page.emulate_media(media="screen")
            if uses_plotly:
                # ensure Plotly actually loaded; if the CDN was slow/blocked, inject it,
                # then wait for every chart to draw its SVG.
                try:
                    await page.wait_for_function("() => !!window.Plotly", timeout=6000)
                except Exception:
                    try:
                        await page.add_script_tag(url="https://cdn.plot.ly/plotly-2.35.2.min.js")
                        await page.wait_for_function("() => !!window.Plotly", timeout=6000)
                    except Exception:
                        pass
                try:
                    await page.wait_for_function(
                        "() => Array.from(document.querySelectorAll('.js-plotly-plot'))"
                        ".every(n => n.querySelector('.main-svg'))",
                        timeout=25000,
                    )
                except Exception:
                    pass
                await page.wait_for_timeout(1500)
            # Wait for the Google Fonts to actually load before capturing — otherwise a slow
            # font fetch means we screenshot fallback fonts (Times/Arial) silently, wrecking
            # the editorial typography the whole design hinges on.
            try:
                await page.evaluate("document.fonts.ready")
            except Exception:
                pass
            # Build the PDF from the SAME per-slide screenshots the preview uses, rather than
            # page.pdf() (which forces print-media emulation and produced blank/black pages).
            # This guarantees PDF == preview == PPTX.
            slide_pngs: List[bytes] = []
            for el in await page.query_selector_all(".slide"):
                slide_pngs.append(await el.screenshot(type="png"))
            await browser.close()
        return {"pdf": _slides_to_pdf(slide_pngs), "slides": slide_pngs}
    finally:
        try:
            os.unlink(tmp.name)
        except Exception:
            pass


def _slides_to_pdf(slide_pngs: List[bytes]) -> bytes:
    """Assemble the full-slide PNG screenshots into a multi-page PDF (one slide per page).
    Uses the same screenshots as the preview/PPTX so all three outputs match exactly."""
    from PIL import Image

    if not slide_pngs:
        raise ValueError("No slides were rendered — cannot build a PDF (the deck has no .slide elements).")
    pages = [Image.open(BytesIO(p)).convert("RGB") for p in slide_pngs]
    buf = BytesIO()
    # Derive DPI from the canvas so the physical page stays 13.333x7.5in at ANY canvas size
    # (a hardcoded 96 DPI assumed a 1280px slide and would emit a 20x11.25in page at 1920px).
    pages[0].save(buf, format="PDF", save_all=True, append_images=pages[1:],
                  resolution=(SLIDE_W_PX / 13.333) * _RENDER_SCALE)
    return buf.getvalue()


def _slides_to_pptx(slide_pngs: List[bytes]) -> bytes:
    """Place each full-slide screenshot onto a 16:9 python-pptx slide."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in slide_pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(BytesIO(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def render_deck(html: str, fmt: str = "pdf") -> bytes:
    """Render the AI's HTML to the requested downloadable format ('pdf' or 'pptx')."""
    fmt = (fmt or "pdf").lower()
    rendered = await _render(html)
    if fmt == "pdf":
        return rendered["pdf"]
    if fmt == "pptx":
        return _slides_to_pptx(rendered["slides"])
    raise ValueError(f"Unsupported format: {fmt!r} (use 'pdf' or 'pptx')")


async def render_slide_images(html: str, *, quality: int = 85, on_progress: ProgressCb = None) -> List[bytes]:
    """Render each slide to a JPEG for in-app preview — same load/Plotly-wait path as
    _render, so the preview matches the downloaded file exactly (charts included)."""
    import os
    import tempfile
    from pathlib import Path
    from playwright.async_api import async_playwright

    if on_progress:
        await on_progress("Rendering charts & slides…")
    tmp = tempfile.NamedTemporaryFile(mode="w", suffix=".html", delete=False, encoding="utf-8")
    try:
        tmp.write(_prepare_html_for_render(html))
        tmp.close()
        uses_plotly = "Plotly.newPlot" in html or "plot.ly" in html
        async with async_playwright() as p:
            browser = await p.chromium.launch(args=["--no-sandbox"])
            page = await browser.new_page(viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX},
                                          device_scale_factor=_RENDER_SCALE)
            await page.goto(Path(tmp.name).as_uri(), wait_until="networkidle")
            await page.emulate_media(media="screen")
            if uses_plotly:
                try:
                    await page.wait_for_function("() => !!window.Plotly", timeout=6000)
                except Exception:
                    try:
                        await page.add_script_tag(url="https://cdn.plot.ly/plotly-2.35.2.min.js")
                        await page.wait_for_function("() => !!window.Plotly", timeout=6000)
                    except Exception:
                        pass
                try:
                    await page.wait_for_function(
                        "() => Array.from(document.querySelectorAll('.js-plotly-plot'))"
                        ".every(n => n.querySelector('.main-svg'))",
                        timeout=25000,
                    )
                except Exception:
                    pass
                await page.wait_for_timeout(1500)
            # Same font-load gate as _render so the preview matches the file exactly.
            try:
                await page.evaluate("document.fonts.ready")
            except Exception:
                pass
            imgs: List[bytes] = []
            for el in await page.query_selector_all(".slide"):
                imgs.append(await el.screenshot(type="jpeg", quality=quality))
            await browser.close()
        return imgs
    finally:
        try:
            os.unlink(tmp.name)
        except Exception:
            pass


async def generate_deck_from_pdf(
    pdf_bytes: bytes,
    *,
    fmt: str = "pdf",
    prompt: Optional[str] = None,
    provider: str = "qwen3.7-max",
    brand: Optional[str] = None,
    structure: Optional[str] = None,
    render: bool = True,
    images: bool = True,
    notes: str = "",
    seed: Optional[str] = None,
    on_progress: ProgressCb = None,
    creativity: str = "balanced",
    pipeline: str = "single",
    models: Optional[Dict[str, str]] = None,
    theme_mode: str = "tbs",
    custom_color: Optional[str] = None,
    style: str = "tbs",
) -> Dict:
    """Full PDF→deck flow: extract the PDF's data, have the AI design the deck with the
    chosen prompt + provider. Renders to the file unless render=False (deferred to download).

    `seed` (e.g. the uploaded file's name) pins a deterministic theme preset per client."""
    from services.pdf_extract import extract_pdf_text
    from services.highlights import to_brief_block

    if on_progress:
        await on_progress("Reading the PDF…")
    data_text = extract_pdf_text(pdf_bytes)
    if not data_text.strip():
        raise ValueError("No extractable text found in the PDF.")
    data_text = data_text + to_brief_block(notes)
    # Resolve the palette (uploaded PDFs have no site to detect, so 'site' collapses to the TBS
    # default; 'custom' honours the picked hex).
    palette = {"accent": TBS_PALETTE["accent"], "accent2": TBS_PALETTE["accent2"]}
    if (theme_mode or "tbs").lower() == "custom" and custom_color:
        from services.site_theme import _accents, _hex_to_rgb
        rgb = _hex_to_rgb(custom_color)
        if rgb:
            palette = _accents(rgb)
    brand_directive = brand or (
        "Design a clean, modern, professional agency-grade presentation. Build the palette around "
        f"{palette['accent']} as --accent and {palette['accent2']} as --accent-2; take the ground, "
        "surface and fonts from the assigned HOUSE STYLE / THEME."
    )
    # Shared cache overlaps image generation with the streamed slide-writing.
    image_cache = {} if images else None
    artifacts = {}   # filled with per-slide md/html by the per-slide pipeline
    html = await generate_deck_html(
        data_text,
        prompt=prompt,
        brand=brand_directive,
        structure=structure or GOOGLE_ADS_STRUCTURE,
        provider=provider,
        on_progress=on_progress,
        image_cache=image_cache,
        seed=seed,
        creativity=creativity,
        pipeline=pipeline,
        models=models,
        style=style,
        artifacts=artifacts,
    )
    html = (await resolve_ai_images(html, on_progress=on_progress, image_cache=image_cache)
            if images else _AI_IMG_RE.sub("", html))
    html = resolve_ai_icons(html)
    html = _apply_theme(html, palette["accent"], palette["accent2"])
    file_bytes = await render_deck(html, fmt=fmt) if render else None
    return {"format": fmt, "html": html, "data_text": data_text, "file_bytes": file_bytes,
            "artifacts": artifacts}
