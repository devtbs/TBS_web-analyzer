#!/usr/bin/env python3
"""
highlight_to_pptx.py — Build a PowerPoint deck from "/on <date>"-highlighted notes.

Scans a body of unstructured notes for lines that begin with the trigger
"/on <date>" and turns each one into its own slide.

GUARDRAIL: slide wording is copied VERBATIM from the matched line. There is no
LLM, paraphrasing, summarising, or rewriting anywhere in this script. Any line
that is not explicitly highlighted with the trigger is ignored completely.

Usage:
    python highlight_to_pptx.py notes.txt -o schedule.pptx
    type notes.txt | python highlight_to_pptx.py -o schedule.pptx
    python highlight_to_pptx.py                 # runs the built-in demo
"""
from __future__ import annotations

import argparse
import sys
from typing import List, Optional

from pptx import Presentation
from pptx.util import Pt

# Shared highlight logic — one source of truth, also used by the AI presentation flow.
from services.highlights import Highlight, extract_highlights


def build_presentation(highlights: List[Highlight], output_path: str) -> None:
    """Write a .pptx with a title slide + one slide per highlight (text inserted verbatim)."""
    prs = Presentation()

    # Cover slide
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Highlights"
    try:
        cover.placeholders[1].text = f"{len(highlights)} highlighted item(s)"
    except (KeyError, IndexError):
        pass

    content_layout = prs.slide_layouts[1]  # "Title and Content"
    for h in highlights:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = h.date          # verbatim
        body = slide.placeholders[1].text_frame
        body.text = h.text                        # verbatim — the full highlighted line
        for paragraph in body.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)

    prs.save(output_path)


DEMO_NOTES = """We had a team sync this morning.
/on 26 may there is a major product launch happening at 9 AM.
I need to remember to buy coffee later.
/on 30 may we have the final client sign-off."""


def _read_input(path: Optional[str]) -> str:
    if path:
        with open(path, "r", encoding="utf-8") as fh:
            return fh.read()
    if not sys.stdin.isatty():          # piped input
        return sys.stdin.read()
    return DEMO_NOTES                    # no file, no pipe -> demo


def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Build a PPTX from /on-highlighted notes.")
    parser.add_argument("input", nargs="?", help="Notes file (omit to read stdin, or run the demo).")
    parser.add_argument("-o", "--output", default="highlights.pptx", help="Output .pptx path.")
    args = parser.parse_args(argv)

    highlights = extract_highlights(_read_input(args.input))
    if not highlights:
        print("No '/on <date>' highlights found — nothing to generate.", file=sys.stderr)
        return 1

    build_presentation(highlights, args.output)
    print(f"Extracted {len(highlights)} highlight(s):")
    for i, h in enumerate(highlights, 1):
        print(f"  {i}. {h.text}")
    print(f"Saved -> {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
